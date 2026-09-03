import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    PRIMARY = RGBColor(234, 88, 12)       # Saffron / Orange (#ea580c)
    PRIMARY_DARK = RGBColor(154, 52, 18)  # Deep Saffron (#9a3412)
    SECONDARY = RGBColor(22, 101, 52)     # Forest Green (#166534)
    ACCENT_TEAL = RGBColor(13, 148, 136)  # Teal (#0d9488)
    TEXT_DARK = RGBColor(30, 41, 59)      # Slate 800 (#1e293b)
    TEXT_MUTED = RGBColor(100, 116, 139)  # Slate 500 (#64748b)
    BG_CREAM = RGBColor(253, 251, 247)    # Light Cream (#fdfbf7)
    CARD_BG = RGBColor(255, 255, 255)     # Pure White
    CARD_BORDER = RGBColor(226, 232, 240) # Slate 200
    HIGHLIGHT_RED = RGBColor(220, 38, 38) # Red (#dc2626)
    HIGHLIGHT_GREEN = RGBColor(21, 128, 61) # Green (#15803d)

    def set_slide_bg(slide, color=BG_CREAM):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=None):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY

        # Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.6))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        # Subtitle if present
        if subtitle:
            s_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.4))
            tf_s = s_box.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
            p_s = tf_s.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(13)
            p_s.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    # =========================================================================
    # SLIDE 1: COVER SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1, RGBColor(255, 248, 232))

    # Decorative Header Band
    top_band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = PRIMARY
    top_band.line.fill.background()

    # Title & Subtitle Card
    c1 = add_card(s1, 1.0, 1.2, 11.333, 5.1, CARD_BG, RGBColor(254, 215, 170))
    tf1 = c1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.8)
    tf1.margin_top = Inches(0.8)
    tf1.margin_right = Inches(0.8)

    p1 = tf1.paragraphs[0]
    p1.text = "THE ART OF LIVING • SRI SRI YOGA"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY

    p2 = tf1.add_paragraph()
    p2.text = "Subscription & Auto-Payment UX Strategy"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_DARK
    p2.space_before = Pt(14)

    p3 = tf1.add_paragraph()
    p3.text = "Optimizing Purchase Conversion, Frictionless Onboarding & Subscription Lifecycle Management"
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(10)

    p4 = tf1.add_paragraph()
    p4.text = "Key Mandates: Low-Friction Purchase Flow • Pay First -> Complete Profile After • 1-Time vs Recurring Economics • Pause/Resume Architecture • Competitor Benchmark"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = PRIMARY_DARK
    p4.space_before = Pt(36)

    # =========================================================================
    # SLIDE 2: EXECUTIVE SUMMARY & DECISION MANDATE
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2)
    add_header(s2, "Executive Mandate", "Executive Summary: Transforming Sri Sri Yoga Subscription UX", "Strategic findings and recommendations for leadership decision-making")

    # 3 Strategic Pillars
    c2_1 = add_card(s2, 0.8, 1.8, 3.6, 5.0)
    tf = c2_1.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "1. FRICTION ELIMINATION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "Move Registration Post-Payment"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "• Current flow asks for 8+ personal fields (Name, Age, PIN, City) BEFORE payment.\n• Creates peak friction right when purchase intent is highest.\n• Recommendation: Adopt 'Pay First -> Complete Profile After' principle to maximize payment conversion."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    c2_2 = add_card(s2, 4.85, 1.8, 3.6, 5.0)
    tf = c2_2.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "2. REVENUE MODEL"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p = tf.add_paragraph()
    p.text = "Evaluate ~20% Renewal Reality"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "• With ~20% current renewal rate, forced auto-debit creates mandate drop-offs and customer anxiety.\n• Competitors (Isha, Habuild) use 1-time upfront or offer transparent auto-renew with clear incentives.\n• Recommendation: Offer 1-time upfront as standard, with optional Auto-Renew discount/perks."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    c2_3 = add_card(s2, 8.9, 1.8, 3.6, 5.0)
    tf = c2_3.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "3. RETENTION & TRUST"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Pause, Cancel & Teacher Profiles"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "• Cult-style Pause Entitlement (15/30/45 days) prevents churn due to travel/illness.\n• Transparent cancellation (no dark patterns) builds high spiritual trust.\n• Teacher profiles & credentials on landing page reduce uncertainty before checkout."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Save presentation
    prs.save("Daily_Yoga_Subscription_UX_Strategy_Updated.pptx")
    print("Base presentation created with 2 slides.")

create_deck()
