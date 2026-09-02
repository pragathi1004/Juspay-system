import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_master_presentation_decks():
    # =========================================================================
    # COLOR PALETTE & DESIGN SYSTEM (Art of Living / Sri Sri Yoga Brand Guidelines)
    # =========================================================================
    PRIMARY = RGBColor(234, 88, 12)        # Saffron / Orange (#ea580c)
    PRIMARY_LIGHT = RGBColor(255, 237, 213)# Light Saffron (#ffedd5)
    PRIMARY_DARK = RGBColor(154, 52, 18)   # Deep Saffron (#9a3412)
    SECONDARY = RGBColor(22, 101, 52)      # Forest Green (#166534)
    SECONDARY_LIGHT = RGBColor(220, 252, 231) # Light Green (#dcfce7)
    ACCENT_BLUE = RGBColor(2, 132, 199)    # Sky Blue (#0284c7)
    ACCENT_BLUE_LIGHT = RGBColor(224, 242, 254)
    TEXT_DARK = RGBColor(15, 23, 42)       # Slate 900 (#0f172a)
    TEXT_MUTED = RGBColor(100, 116, 139)   # Slate 500 (#64748b)
    BG_CREAM = RGBColor(253, 251, 247)     # Warm Cream (#fdfbf7)
    BG_WHITE = RGBColor(255, 255, 255)     # Pure White
    CARD_BORDER = RGBColor(226, 232, 240)  # Slate 200
    CARD_ORANGE_BORDER = RGBColor(254, 215, 170)
    ALERT_RED = RGBColor(220, 38, 38)      # Red (#dc2626)
    ALERT_RED_BG = RGBColor(254, 242, 242) # Light Red (#fef2f2)
    GOLD = RGBColor(217, 119, 6)

    # =========================================================================
    # HELPER FUNCTIONS
    # =========================================================================
    def set_bg(slide, prs, color=BG_CREAM):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=None):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY

        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.55))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        if subtitle:
            s_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.28), Inches(11.7), Inches(0.35))
            tf_s = s_box.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
            p_s = tf_s.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(11.5)
            p_s.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, bg_color=BG_WHITE, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    # =========================================================================
    # DECK 1: 36-SLIDE BUSINESS, UX STRATEGY & PRODUCT PRESENTATION
    # =========================================================================
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Slide 1: Cover Slide
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, prs, RGBColor(255, 248, 232))
    top_band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = PRIMARY
    top_band.line.fill.background()
    c1 = add_card(s1, 0.8, 1.1, 11.733, 5.3, BG_WHITE, CARD_ORANGE_BORDER)
    tf = c1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.6)
    p = tf.paragraphs[0]
    p.text = "THE ART OF LIVING • SRI SRI YOGA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "Subscription Journey, Payment UX &\nLifecycle Strategy"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)
    p = tf.add_paragraph()
    p.text = "Comprehensive Research, Competitor Benchmarks, 'Pay First -> Profile After' Architecture & Executive Roadmap"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "Prepared for: Product, Business, Legal & Leadership Teams | September 2026"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p.space_before = Pt(32)

    # Slide 2: Executive Summary & Core Verdict
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, prs)
    add_header(s2, "Executive Briefing", "Executive Summary: Transforming Sri Sri Yoga Digital Subscriptions", "Strategic recommendations to maximize checkout conversion and curb subscriber churn")
    cards_s2 = [
        ("1. Pay First Architecture", "Shift from pre-payment registration to 'Payment First -> Progressive Profile After'. Eliminates high pre-checkout drop-off.", PRIMARY, PRIMARY_LIGHT),
        ("2. 3-Question Onboarding", "Streamline pre-pay survey to 3 engaging questions (<10s). Remove 'Time dedication' question due to fixed class timings.", ACCENT_BLUE, ACCENT_BLUE_LIGHT),
        ("3. Hybrid Payment Model", "Default to 1-Time upfront purchase (3M/6M/12M) with optional Auto-Renew perk to navigate ~20% baseline renewal rate.", SECONDARY, SECONDARY_LIGHT),
        ("4. Lifecycle Safeguards", "Implement structured Pause pools (15d/30d/45d) and transparent self-serve cancellation to establish high brand trust.", GOLD, RGBColor(254, 243, 199))
    ]
    for i, (title, desc, col, bg_col) in enumerate(cards_s2):
        c = add_card(s2, 0.8 + i*2.98, 1.75, 2.8, 5.1, bg_col, CARD_BORDER)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(12)

    # Slide 3: Core Business Problem Diagnosis
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, prs)
    add_header(s3, "Problem Statement", "The Core Business Problem: High Pre-Payment Friction Peak", "Current registration form demands extensive personal data before value delivery or purchase commitment")
    c3_l = add_card(s3, 0.8, 1.75, 5.7, 5.1, ALERT_RED_BG, RGBColor(254, 202, 202))
    tf = c3_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "🔴 CURRENT FRICTION-HEAVY JOURNEY"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "Landing Page -> Select Plan -> OTP Login -> 4 Questions -> 8-Field Registration Form (Name, Age, WhatsApp, PIN, City) -> Payment Gateway"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)
    p = tf.add_paragraph()
    p.text = "• Friction Peak: Asking for age, postal code and full personal details BEFORE payment creates 40-55% drop-off at peak purchase intent.\n• Psychological Barrier: Users question why an online yoga course requires postal address before payment.\n• Mobile Typing Fatigue: 8 manual text fields on mobile leads to cart abandonment."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    c3_r = add_card(s3, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c3_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "🟢 RECOMMENDED CONVERSION-FIRST PRINCIPLE"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Landing Page -> Select Plan -> OTP (Mobile/WhatsApp) -> 3 Quick 1-Tap Questions -> PAYMENT -> Success -> Progressive Profile Setup"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)
    p = tf.add_paragraph()
    p.text = "• Value First: Move financial transaction forward while buyer intent is highest.\n• 100% Commitment: Once paid, users gladly complete profile details to join live batches.\n• Skip-for-now option: Allows immediate access to live yoga classes without blocking."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    # Slide 4: Flow A (Recommended: Payment First)
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, prs)
    add_header(s4, "User Flow Architecture", "Flow A (Recommended): Payment First Journey Architecture", "Optimal flow with minimal pre-payment survey and post-payment progressive profile setup")
    steps_a = [
        ("Step 1: Landing Page", "Explores benefits, schedule & teacher profiles. Clicks 'Try for Free' -> Smooth scrolls to pricing."),
        ("Step 2: Plan Selection", "Chooses 3M / 6M / 12M membership tier with clear per-month anchoring."),
        ("Step 3: Lightweight OTP", "Quick 4-digit OTP via Mobile SMS or WhatsApp. Auto-populates phone number."),
        ("Step 4: 3-Question Survey", "1-tap selection: Intent, Experience, and Desired Transformation. Takes <10 seconds."),
        ("Step 5: Juspay Checkout", "Instant payment via UPI, Cards, NetBanking or optional AutoPay mandate."),
        ("Step 6: Post-Pay Profile", "Personalizes yoga batch, health notes, age & city with 'Skip for now' option.")
    ]
    for i, (title, desc) in enumerate(steps_a):
        c = add_card(s4, 0.8 + (i%3)*3.95, 1.75 + (i//3)*2.6, 3.8, 2.35, BG_WHITE, CARD_ORANGE_BORDER if i==4 else CARD_BORDER)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK if i==4 else TEXT_DARK
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(6)

    # Slide 5: Pre-Payment Questionnaire Optimization
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, prs)
    add_header(s5, "Survey Design", "Pre-Payment Questionnaire: 3 High-Value, Low-Friction Questions", "Optimizing data collection: engaging the user without creating checkout drop-offs")
    q_cards = [
        ("Question 1: User Intent", "What brings you to Sri Sri Yoga?", "• Improve flexibility & mobility\n• Build strength & stamina\n• Relieve stress & anxiety\n• Weight management\n• Relieve chronic body pain\n• Mindfulness & meditation", PRIMARY),
        ("Question 2: Starting Level", "What is your yoga experience?", "• Complete Beginner (Never practiced)\n• Some Experience (Occasional practice)\n• Regular Practitioner (1-2x per week)\n• Advanced Yogi (Daily practitioner)", ACCENT_BLUE),
        ("Question 3: Desired Goal (NEW)", "What is your primary wellness goal?", "• Establishing a daily healthy habit\n• Higher daily energy & vitality\n• Better sleep quality & mental calm\n• Guided posture & breathing mastery", SECONDARY)
    ]
    for i, (q_t, q_sub, q_opts, col) in enumerate(q_cards):
        c = add_card(s5, 0.8 + i*3.95, 1.75, 3.8, 4.0, BG_WHITE, CARD_BORDER)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
        p = tf.paragraphs[0]
        p.text = q_t
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = f'"{q_sub}"'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(4)
        p = tf.add_paragraph()
        p.text = q_opts
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    c5_b = add_card(s5, 0.8, 5.95, 11.733, 1.0, ALERT_RED_BG, RGBColor(254, 202, 202))
    tf = c5_b.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = "⚠️ REMOVED QUESTION: 'How much time can you dedicate to yoga each day?'"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "Rationale: Sri Sri Yoga sessions run on fixed live schedule slots (6:00 AM, 7:30 AM, 5:30 PM, 6:30 PM). Asking time dedication is non-actionable and creates unnecessary cognitive friction."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(2)

    # Slide 6: Flow B (Alternate: Minimal Pre-Payment Details)
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, prs)
    add_header(s6, "Alternate Flow", "Flow B (Alternate): Minimal Details Collected Pre-Payment", "Evaluating a flow where only basic name and WhatsApp are captured prior to payment gateway")
    c6_l = add_card(s6, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c6_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "FLOW B ARCHITECTURE:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p = tf.add_paragraph()
    p.text = "1. Landing Page -> Select Plan\n2. OTP Authentication (Mobile)\n3. 2 Quick Survey Questions\n4. Minimal Pre-Pay Form (Full Name & Email only)\n5. Juspay Payment Gateway\n6. Subscription Activation\n7. Post-Payment: Remaining Health & Address Profile"
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "• Intended Benefit: Captures lead name for abandoned cart remarketing.\n• Operational Cost: Adds 1 additional step before payment compared to Flow A."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(10)

    c6_r = add_card(s6, 6.8, 1.75, 5.7, 5.1, RGBColor(255, 251, 235), GOLD)
    tf = c6_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "WHEN IS FLOW B VIABLE?"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p = tf.add_paragraph()
    p.text = "• Abandoned Cart Recovery: If marketing relies heavily on WhatsApp outreach for users who drop off at checkout.\n• Tax Invoice Legalities: If instant GST invoice generation requires full legal name upfront.\n• A/B Testing Recommendation: Launch Flow A as 80% primary traffic, run Flow B as 20% challenger to validate whether pre-payment name collection reduces overall payment conversion."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 7: Side-by-Side Comparison: Flow A vs Flow B
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, prs)
    add_header(s7, "Comparative Evaluation", "Side-by-Side Comparison: Flow A (Pay First) vs Flow B (Minimal Pre-Pay)", "Direct trade-off analysis across conversion, user friction, CRM integration and legal considerations")
    
    # Table comparison
    table_shape = s7.shapes.add_table(7, 3, Inches(0.8), Inches(1.75), Inches(11.733), Inches(4.9))
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.733)
    tbl.columns[1].width = Inches(4.5)
    tbl.columns[2].width = Inches(4.5)

    headers = ["Evaluation Metric", "Flow A: Payment First (Recommended)", "Flow B: Minimal Details Pre-Payment"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_DARK if j>0 else TEXT_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = BG_WHITE

    matrix = [
        ("Checkout Conversion Rate", "Highest (Estimated +25% to +40% uplift)", "Moderate (+10% to +15% uplift vs current)"),
        ("Pre-Payment Time to Checkout", "< 30 seconds total", "45 - 60 seconds"),
        ("Pre-Payment Friction", "Minimal (OTP + 3 one-tap taps)", "Medium (typing Name & Email on mobile)"),
        ("CRM Lead Capture for Drop-offs", "Captured via Mobile OTP (Phone number captured)", "Captured via Name + Phone + Email"),
        ("Post-Payment Profile Completion", "85-92% expected completion (motivated post-pay)", "90-95% expected completion"),
        ("Final Strategic Recommendation", "Adopt Flow A as Primary Default", "Reserve as A/B Testing Variant")
    ]
    for i, row in enumerate(matrix):
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if i%2==0 else RGBColor(248, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(9.5)
            p.font.bold = (j==0 or (j==1 and i==5))
            p.font.color.rgb = PRIMARY if (j==1 and i==5) else TEXT_DARK

    # Slide 8: Before vs After Visual Transformation
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8, prs)
    add_header(s8, "Transformation", "Before vs After: Journey Streamlining & Friction Elimination", "Visual comparison of current high-friction user journey vs recommended frictionless flow")
    c8_l = add_card(s8, 0.8, 1.75, 5.7, 5.1, ALERT_RED_BG, RGBColor(254, 202, 202))
    tf = c8_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "❌ CURRENT FLOW (BEFORE): 8 STEPS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "1. Landing Page -> 2. Select Plan -> 3. OTP Login -> 4. Intent Question -> 5. Experience Question -> 6. Time Dedication Question -> 7. HEAVY 8-FIELD REGISTRATION FORM (First/Last Name, Email, WhatsApp, Age, PIN, City) -> 8. Payment Gateway"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "💥 Result: Major drop-off at Step 7. User leaves before completing payment."
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p.space_before = Pt(12)

    c8_r = add_card(s8, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c8_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "✅ RECOMMENDED FLOW (AFTER): 5 STEPS TO PAYMENT"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "1. Landing Page -> 2. Select Plan -> 3. Lightweight OTP -> 4. 3 Quick Survey Taps (<10s) -> 5. JUSPAY PAYMENT GATEWAY -> 6. Instant Active Subscription -> 7. Progressive Profile Setup"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "🚀 Result: Frictionless path to payment. 100% of paying users enter the retention funnel."
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p.space_before = Pt(12)

    # Slide 9: Post-Payment Progressive Profile Setup
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9, prs)
    add_header(s9, "Profile Experience", "Post-Payment Experience: 'Complete Your Profile' Screen", "Framing profile collection as an empowering personalization step rather than a registration barrier")
    c9_l = add_card(s9, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c9_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "POST-PAYMENT PSYCHOLOGY & UX:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Framing: 'Your Sri Sri Yoga Subscription is Active! Personalize your daily practice...'\n• Progressive Disclosure: Present fields in intuitive logical blocks (Personal Info -> Practice Batch -> Health Notes).\n• Trust Multiplier: Since payment is already confirmed, user views profile questions as teacher preparation rather than data harvesting.\n• Skip Option: 'Skip for now & start practicing' ensures zero blockage to live Zoom sessions."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c9_r = add_card(s9, 6.8, 1.75, 5.7, 5.1, RGBColor(240, 249, 255), ACCENT_BLUE)
    tf = c9_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "FIELDS COLLECTED POST-PAYMENT:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "1. Full Name (For official Sri Sri Yoga certificate)\n2. Email Address (For calendar invites & receipts)\n3. Age & Gender (For customized posture variations)\n4. City & State (For regional ashram events)\n5. Preferred Batch Timing (6:00 AM / 7:30 AM / 5:30 PM / 6:30 PM)\n6. Health & Injury Notes (Back pain, knee issues, blood pressure - shared confidentially with certified teachers)"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 10: Landing Page Trust & Teacher Profiles
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10, prs)
    add_header(s10, "Landing Page Architecture", "Landing Page Trust: Certified Teacher Profiles ('Meet Your Faculty')", "Establishing spiritual lineage, authenticity, and instructor credibility prior to purchase decision")
    c10_l = add_card(s10, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c10_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "WHY TEACHER PROFILES ARE ESSENTIAL:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Spiritual & Pedagogical Trust: Art of Living students value direct teacher connection and certified ashram lineage.\n• Expertise Transparency: Showcasing 10+ years of teaching experience, AYUSH ministry certifications, and international ashram training builds confidence.\n• Human Connection: Converts a faceless subscription into a personalized mentorship with trusted yoga masters.\n• Reduced Purchase Anxiety: First-time yogis feel reassured that qualified teachers will guide posture alignments."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c10_r = add_card(s10, 6.8, 1.75, 5.7, 5.1, RGBColor(255, 247, 237), CARD_ORANGE_BORDER)
    tf = c10_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "RECOMMENDED FACULTY CARD SPECIFICATION:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• High-Quality Studio Photo: Warm, serene, professional portrait in traditional yoga attire.\n• Name & Title: e.g., 'Swami Paramjyoti' / 'Priya Sharma, Senior Sri Sri Yoga Faculty'.\n• Ashram Lineage & Experience: e.g., '14+ Years Experience | Trained at Bangalore International Ashram'.\n• Specialization Badges: Pranayama, Hatha Yoga, Sukshma Vyayama, Therapeutic Alignment.\n• Faculty Quote: A short 1-line inspiring philosophy from the teacher."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 11: Competitor Benchmark Matrix (6 Competitors)
    s11 = prs.slides.add_slide(blank_layout)
    set_bg(s11, prs)
    add_header(s11, "Market Intelligence", "Competitive Landscape: Benchmark of 6 Leading Indian Wellness Brands", "Comparative analysis of Times Health+, Kamya Health, Isha Yoga, Satvic Movement, Cult.fit, and Habuild")
    
    t11 = s11.shapes.add_table(7, 5, Inches(0.8), Inches(1.75), Inches(11.733), Inches(4.9))
    tbl11 = t11.table
    tbl11.columns[0].width = Inches(2.2)
    tbl11.columns[1].width = Inches(2.4)
    tbl11.columns[2].width = Inches(2.2)
    tbl11.columns[3].width = Inches(2.5)
    tbl11.columns[4].width = Inches(2.433)

    c_hdrs = ["Platform", "Pricing & Duration", "Payment Model", "Registration Timing", "Pause / Cancellation"]
    for j, h in enumerate(c_hdrs):
        cell = tbl11.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_DARK if j>0 else TEXT_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = BG_WHITE

    c_rows = [
        ("Habuild Yoga", "₹1,999 / 3 Mo (WhatsApp led)", "1-Time Upfront", "Payment first, profile on WhatsApp", "WhatsApp support pause"),
        ("Cult / Cult.fit", "₹14,990 / yr (Cultpass)", "1-Time + Auto-Debit", "Account login -> Pay -> Profile", "App pause slider (15-45 days)"),
        ("Satvic Movement", "$40/mo, $90/3mo, $290/yr", "Auto-Renew Default", "Minimal checkout form -> Pay", "Self-serve customer portal"),
        ("Isha Yoga / Inner Eng.", "₹1,500 - ₹3,500 / program", "1-Time Enrollment", "Enrollment fee -> Portal profile", "Non-refundable, transferrable"),
        ("Times Health+", "₹999 / mo, ₹5,999 / yr", "Recurring Auto-Debit", "Account login -> Pay -> Schedule", "Self-serve cancel in 1-click"),
        ("Kamya Health", "₹1,499 / 21-Day Challenge", "1-Time Challenge Fee", "Challenge sign-up -> Community", "Non-refundable fixed cohort")
    ]
    for i, row in enumerate(c_rows):
        for j, val in enumerate(row):
            cell = tbl11.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if i%2==0 else RGBColor(248, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(8.5)
            p.font.bold = (j==0)
            p.font.color.rgb = TEXT_DARK

    # Slide 12: Cult.fit Deep Dive & Package Mechanics
    s12 = prs.slides.add_slide(blank_layout)
    set_bg(s12, prs)
    add_header(s12, "Competitor Deep Dive", "Cult.fit Package Architecture & Subscription Mechanics", "Analyzing Cultpass subscription tiers, pause entitlement slider, and cancellation UX")
    c12_l = add_card(s12, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c12_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "CULT.FIT SUBSCRIPTION MECHANICS:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Tiered Package Structure: Cultpass ELITE (All centers + At-Home), Cultpass PRO (Gyms), Cultpass HOME (Digital online classes).\n• Pause Entitlement Pool: 3-Month Plan (15 pause days), 6-Month Plan (30 pause days), 12-Month Plan (45 pause days).\n• Seamless In-App Pause: Interactive slider allows members to select start date and end date. Automatically recalculates subscription expiry date.\n• Early Unpause Credit: Resuming early refunds unused pause days back to the user's available pool immediately."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c12_r = add_card(s12, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c12_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "WHAT SRI SRI YOGA SHOULD ADOPT:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "✅ 15d / 30d / 45d Pause Entitlement Matrix: Proven standard in Indian fitness/wellness industry.\n✅ Early Resume Refund: Prevents user frustration if travel plans change.\n✅ Clear Visual Expiry Display: Showing 'Your subscription will now expire on [New Date]' builds trust.\n❌ Avoid Cult's Heavy Transfer Fees: Keep pause and rescheduling 100% free for Sri Sri Yoga members."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 13: Habuild Deep Dive (WhatsApp-First Simplicity)
    s13 = prs.slides.add_slide(blank_layout)
    set_bg(s13, prs)
    add_header(s13, "Competitor Deep Dive", "Habuild Deep Dive: WhatsApp-Led Onboarding & Retention", "How Habuild achieved scale via zero-friction checkout and WhatsApp daily habit formation")
    c13_l = add_card(s13, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c13_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "HABUILD MODEL HIGHLIGHTS:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• WhatsApp-First Ecosystem: Users sign up with only Mobile number -> Instant payment -> All daily Zoom links and reminders delivered via WhatsApp bot.\n• Upfront 1-Time Purchases: ₹1,999 for 3 Months or ₹3,999 for Annual pass. No forced auto-debit resistance.\n• Attendance Streak Gamification: WhatsApp bot tracks daily attendance streaks, badges, and teacher shout-outs.\n• Minimal Web UI: Zero complicated web portals; all customer interactions live where users already spend time."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c13_r = add_card(s13, 6.8, 1.75, 5.7, 5.1, RGBColor(240, 249, 255), ACCENT_BLUE)
    tf = c13_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "KEY TAKEAWAYS FOR ART OF LIVING:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "1. WhatsApp OTP & Link Delivery: Integrate Juspay + WhatsApp notifications for instant class joining links.\n2. Daily Habit Building: Send daily morning reminders 15 minutes before chosen batch.\n3. Frictionless Upfront Pricing: Offer 1-time upfront payments to avoid mandate drop-offs among older demographics."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 14: One-Time Upfront Payment vs Recurring Auto-Payment
    s14 = prs.slides.add_slide(blank_layout)
    set_bg(s14, prs)
    add_header(s14, "Payment Economics", "Payment Model Trade-Offs: 1-Time Upfront vs Recurring Auto-Payment", "Evaluating customer friction, mandate resistance, revenue predictability, and chargeback risks")
    c14_l = add_card(s14, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c14_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "1-TIME UPFRONT PAYMENT (3M / 6M / 12M)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "✅ PROS:\n• Zero Mandate Friction: No e-mandate registration or card tokenization failures.\n• High Conversion: Indian consumers readily pay upfront for fixed-term wellness.\n• Zero Chargeback/Dispute Risk: No surprise debits or customer complaints.\n• Complete Trust: Avoids psychological fear of 'hidden recurring subscriptions'.\n\n❌ CONS:\n• Requires active renewal campaigns at cycle end.\n• Unpredictable revenue if reactivation cadence is weak."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)

    c14_r = add_card(s14, 6.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c14_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "RECURRING AUTO-PAYMENT (AUTOPAY / SI)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "✅ PROS:\n• High Lifetime Value (LTV): Automatic renewal maximizes long-term retention.\n• Predictable Cash Flow: Consistent monthly/annual recurring revenue (ARR).\n• Reduced Operational Overhead: System handles recurring collection automatically.\n\n❌ CONS:\n• Mandate Drop-Off: Upfront e-mandate authentication causes 30-40% drop-off.\n• Customer Hesitation: Users resist recurring debits for habit-based products.\n• Compliance Overhead: Strict RBI 24h pre-debit alerts & cancellation rules required."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)

    # Slide 15: The ~20% Renewal Rate Reality & Retention Economics
    s15 = prs.slides.add_slide(blank_layout)
    set_bg(s15, prs)
    add_header(s15, "Renewal Reality", "The ~20% Renewal Rate Reality & Retention Economics", "Why forced auto-debit fails when organic renewal is ~20%, and the strategic metrics we must track")
    c15_l = add_card(s15, 0.8, 1.75, 5.7, 5.1, ALERT_RED_BG, RGBColor(254, 202, 202))
    tf = c15_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "📊 BUSINESS REALITY: ~20% RENEWAL BASELINE"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "• Strategic Risk: If only ~20% of users naturally renew, forcing 100% of new users into recurring auto-debit triggers severe customer resistance, high cancellation rates, and payment gateway chargeback disputes.\n• Root Causes of Low Renewal:\n  1. Lack of daily attendance habit formation.\n  2. Absence of mid-cycle motivation or teacher check-ins.\n  3. Friction when users travel or fall sick without pause options.\n  4. Weak renewal communication (sent only on the day of expiry)."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c15_r = add_card(s15, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c15_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "📈 MANDATORY LIFECYCLE METRICS TO TRACK:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "1. Activation Rate: % of paid users attending First Live Session within 48h.\n2. WAU / Habit Index: Weekly attendance frequency (Target: >= 3 sessions/wk).\n3. 30-Day & 45-Day Retention: % active at end of Month 1.\n4. Pause Utilization: % using pause feature vs churn.\n5. Renewal Communication Conversion: Efficacy of T-14, T-7, and T-1 WhatsApp reminders.\n6. Reactivation Rate: % of expired members repurchasing within 60 days."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 16: The "No Renewal" Expiry Model Analysis
    s16 = prs.slides.add_slide(blank_layout)
    set_bg(s16, prs)
    add_header(s16, "Business Models", "Evaluating the 'No Renewal' Model (Fixed Expiry & Voluntary Repurchase)", "Assessing fixed-term cohort purchasing vs perpetual subscription commitments")
    c16_l = add_card(s16, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c16_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "HOW THE 'NO RENEWAL' MODEL WORKS:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Principle: User purchases a defined package (e.g., 3-Month, 6-Month, or 12-Month access). At the end of the period, access expires cleanly unless the user actively chooses to repurchase.\n• Alignment with Indian Wellness: Follows traditional ashram/course enrollment mindset (like Happiness Program or Sahaj Samadhi).\n• Zero Involuntary Churn: No failed billing attempts, grace period headaches, or card expiry issues."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c16_r = add_card(s16, 6.8, 1.75, 5.7, 5.1, RGBColor(255, 251, 235), GOLD)
    tf = c16_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "RECOMMENDED HYBRID MODEL (OPTION C):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p = tf.add_paragraph()
    p.text = "• Default Setting: 1-Time Upfront Purchase (Maximizes initial checkout conversion).\n• Optional Toggle: 'Enable Auto-Renewal & Save 10% / Get 15 Extra Pause Days'.\n• Result: Best of both worlds—captures high-intent recurring subscribers without repelling users who prefer fixed-term upfront payment."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 17: Recommended Pricing Architecture
    s17 = prs.slides.add_slide(blank_layout)
    set_bg(s17, prs)
    add_header(s17, "Pricing Strategy", "Sri Sri Yoga Pricing Architecture: 3M, 6M & 12M Tiers", "Balanced tier structure anchored on monthly value with built-in pause entitlements")
    p_cards = [
        ("3 Months Starter Plan", "₹1,499 / 3 Months", "₹499 / Month equivalent", "• Live Daily Guided Yoga\n• Morning & Evening Batches\n• 15 Days Pause Entitlement\n• Sri Sri Yoga Certificate\n• Web + Mobile App Access", CARD_BORDER),
        ("6 Months Transformation", "₹2,699 / 6 Months", "₹449 / Month (Save 10%)", "• All 3-Month Plan Inclusions\n• 30 Days Pause Entitlement\n• Monthly Masterclasses\n• Posture Correction Workshops\n• Priority Teacher Q&A", CARD_BORDER),
        ("12 Months Annual Yogi ⭐", "₹4,999 / Year", "₹416 / Month (Best Value)", "• Complete 365-Day Access\n• 45 Days Pause Entitlement\n• Unlimited Masterclass Library\n• Ayurvedic Health Consult\n• Dedicated Mentor Support", CARD_ORANGE_BORDER)
    ]
    for i, (p_t, p_pr, p_eq, p_feat, p_bor) in enumerate(p_cards):
        c = add_card(s17, 0.8 + i*3.95, 1.75, 3.8, 5.1, RGBColor(255, 247, 237) if i==2 else BG_WHITE, p_bor)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
        p = tf.paragraphs[0]
        p.text = p_t
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK if i==2 else TEXT_DARK
        p = tf.add_paragraph()
        p.text = p_pr
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY if i==2 else TEXT_DARK
        p.space_before = Pt(4)
        p = tf.add_paragraph()
        p.text = p_eq
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = SECONDARY
        p.space_before = Pt(2)
        p = tf.add_paragraph()
        p.text = p_feat
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(12)

    # Slide 18: Pause & Resume Entitlement Framework
    s18 = prs.slides.add_slide(blank_layout)
    set_bg(s18, prs)
    add_header(s18, "Lifecycle UX", "Pause / Resume Architecture & Entitlement Pool (15d / 30d / 45d)", "Preventing churn during travel, illness, or festivals with self-serve pause and early resume refunds")
    c18_l = add_card(s18, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c18_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "PAUSE ENTITLEMENT RULES BY PLAN:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• 3-Month Plan: 15 Days Pause Pool\n• 6-Month Plan: 30 Days Pause Pool\n• 12-Month Plan: 45 Days Pause Pool\n\nOPERATIONAL RULES:\n1. Minimum Pause Duration: 3 consecutive days.\n2. Maximum Pauses Allowed: Up to 3 pause events per subscription cycle.\n3. Date Extension Math: New Expiry Date = Original Expiry Date + Number of Days Paused.\n4. Expiration: Unused pause days expire when subscription ends (non-transferable)."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c18_r = add_card(s18, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c18_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "EARLY RESUME REFUND MECHANICS:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "• Scenario: User pauses for 10 days, but decides to resume after 4 days.\n• Immediate Recalculation:\n  - Days Consumed: 4 Days.\n  - Days Refunded to Balance: 6 Days.\n  - Subscription Extension: Adjusted by exact 4 days.\n• Class Booking Intercept: If a paused user clicks 'Join Today's Class', a friendly modal prompts: 'Resume & Book Class now?'. Seamless 1-tap unpause."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 19: Cancellation Best Practices & Anti-Dark Patterns
    s19 = prs.slides.add_slide(blank_layout)
    set_bg(s19, prs)
    add_header(s19, "Cancellation UX", "Cancellation Best Practices & Anti-Dark Pattern Architecture", "Transparent self-serve cancellation that protects brand reputation and avoids chargeback disputes")
    c19_l = add_card(s19, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c19_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "ANTI-DARK PATTERN CANCELLATION PRINCIPLES:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Easy to Find: Accessible directly from 'Manage Subscription' in customer profile. No requirement to call customer care or send emails.\n• Transparent Offer of Pause vs Cancel: Present 'Would you prefer to pause for 15 days instead?' without blocking cancellation.\n• Retain Paid Access: User retains full class access until the end of the currently paid billing period.\n• Instant Confirmation: Display clear confirmation on screen + dispatch instant confirmation via WhatsApp and Email."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c19_r = add_card(s19, 6.8, 1.75, 5.7, 5.1, RGBColor(255, 241, 242), ALERT_RED)
    tf = c19_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "DARK PATTERNS TO ELIMINATE:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "❌ Hidden Cancel Buttons in deep submenus.\n❌ Requiring Phone Call to Support Agent to cancel.\n❌ Immediate Access Revocation when user has already paid for remainder of cycle.\n❌ Multi-Step 'Confirm 5 times' Guilt Trips.\n\nEthical cancellation builds long-term spiritual trust, leading to 3x higher voluntary reactivation."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 20: Plan Upgrade / Downgrade Flow
    s20 = prs.slides.add_slide(blank_layout)
    set_bg(s20, prs)
    add_header(s20, "Plan Switching", "Plan Upgrade, Downgrade & Tier Switching Architecture", "Immediate pro-rated upgrades and cycle-end downgrades with clear financial transparency")
    c20_l = add_card(s20, 0.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c20_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "PLAN UPGRADES (e.g., 3M -> 12M ANNUAL):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "• Immediate Activation: Upgraded benefits (masterclasses, extra pause days) activate instantly upon confirmation.\n• Proration Credit: System calculates unused value of current plan and deducts it from the new plan charge.\n• Example Math:\n  - Current 3M Plan: ₹1,499 (45 days remaining = ₹750 credit)\n  - New 12M Plan: ₹4,999\n  - Amount Payable Today: ₹4,999 - ₹750 = ₹4,249.\n• Single Checkout: Processed in 1 click via Juspay."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c20_r = add_card(s20, 6.8, 1.75, 5.7, 5.1, RGBColor(240, 249, 255), ACCENT_BLUE)
    tf = c20_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "PLAN DOWNGRADES (e.g., 12M -> 3M):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "• Effective Date: Takes effect at the end of the current paid billing cycle (no mid-cycle penalty).\n• Continued Access: User retains full annual privileges until the scheduled cycle end date.\n• Mandate Adjustment: Juspay recurring mandate is automatically updated to the lower amount for next cycle.\n• Confirmation: Instant WhatsApp notification confirming scheduled downgrade date."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 21: Free Trial Strategy & Abuse Prevention
    s21 = prs.slides.add_slide(blank_layout)
    set_bg(s21, prs)
    add_header(s21, "Acquisition Strategy", "Free Trial Architecture: 14-Day Evaluation & Abuse Prevention", "Structured free trial rules: 1 trial per user, automated transition notifications, and zero lock-in")
    c21_l = add_card(s21, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c21_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "FREE TRIAL ELIGIBILITY & DURATION:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Duration: 14 Days Full Access (Sufficient to experience 8-10 live yoga sessions and build a morning habit).\n• Eligibility: Strictly 1 Free Trial per unique Mobile Number and Email ID.\n• Zero Upfront Charge: ₹0 charged during trial period.\n• Mandate Registration: Optional ₹1 authorization via Juspay UPI/Card with instant refund for trial auto-conversion."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c21_r = add_card(s21, 6.8, 1.75, 5.7, 5.1, RGBColor(255, 251, 235), GOLD)
    tf = c21_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "TRIAL TRANSITION CADENCE & ALERTS:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p = tf.add_paragraph()
    p.text = "• Day 0: Instant welcome WhatsApp with batch timing and Zoom link.\n• Day 7 (Mid-Trial): Milestone celebration: 'You have completed 5 yoga sessions! Meet your teacher for Q&A.'\n• Day 11 (T-3 Advance Alert): Mandatory WhatsApp notification: 'Your 14-day trial ends in 3 days. Your subscription will activate on [Date]. Cancel anytime for ₹0.'\n• Day 14 (Execution): Trial converts to active paid membership."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 22: 45-Day Access & Multi-Platform Entitlement
    s22 = prs.slides.add_slide(blank_layout)
    set_bg(s22, prs)
    add_header(s22, "Entitlement Model", "45-Day Companion Access & Multi-Platform Entitlement Architecture", "Seamless cross-platform authentication across Web and Art of Living Mobile App")
    c22_l = add_card(s22, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c22_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "45-DAY ACCESS ENTITLEMENT RESOLUTION:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Context: Certain Art of Living flagship programs include 45-day guided yoga companion access.\n• Entitlement Creation: Upon course purchase, system generates a 45-day entitlement token linked to customer's phone number.\n• Expiry Countdown Display: Customer dashboard prominently displays '45-Day Immersion Access: 32 Days Remaining'.\n• Seamless Upgrade Bridge: At Day 35 (T-10), offer an exclusive 'Loyalty Transition Offer' to convert to the Annual 12-Month plan with 20% discount."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c22_r = add_card(s22, 6.8, 1.75, 5.7, 5.1, RGBColor(240, 249, 255), ACCENT_BLUE)
    tf = c22_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "CROSS-PLATFORM APP ENTITLEMENT (WEB + AOL APP):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "• Single Sign-On (SSO): Customer's verified Mobile Number unlocks:\n  1. Web Portal (Daily live Zoom join links, masterclass replays, billing).\n  2. Art of Living Mobile App (Audio meditations, offline sadhana player).\n• Cloud Sync: Attendance streak and profile data sync automatically across Web and Mobile App."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 23: Authentication & Login Cleanup (Apple Removed)
    s23 = prs.slides.add_slide(blank_layout)
    set_bg(s23, prs)
    add_header(s23, "Authentication UX", "Authentication Cleanup: Removing Apple Login & Streamlining Auth", "Optimizing authentication for Indian demographics: Mobile/WhatsApp OTP & Google 1-Tap")
    c23_l = add_card(s23, 0.8, 1.75, 5.7, 5.1, ALERT_RED_BG, RGBColor(254, 202, 202))
    tf = c23_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "❌ REMOVED: APPLE LOGIN OPTION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "• Rationale for Removal:\n  1. Hidden Email Issues: Apple Private Relay generates masked @privaterelay.appleid.com emails, breaking WhatsApp class reminders and CRM record linking.\n  2. Demographics: 88%+ of target Indian yoga audience uses Android / Mobile OTP.\n  3. Friction: Removing Apple Login simplifies the auth modal to a single clean primary action."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c23_r = add_card(s23, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c23_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "✅ RECOMMENDED PRIMARY LOGIN OPTIONS:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "1. Mobile + WhatsApp OTP (Primary):\n   - User enters 10-digit number.\n   - Instant 4-digit OTP via SMS or WhatsApp.\n   - Auto-read OTP on mobile devices.\n\n2. Google 1-Tap (Secondary):\n   - 1-tap fast login on desktop Chrome and Android.\n   - Captures verified email address."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 24: Explicit User Consent & Legal Transparency UI
    s24 = prs.slides.add_slide(blank_layout)
    set_bg(s24, prs)
    add_header(s24, "Compliance UX", "Explicit User Consent & Payment Authorization Architecture", "Ensuring clear, unambiguous consent without deceptive pre-selected checkboxes")
    c24_l = add_card(s24, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c24_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "EXPLICIT CONSENT STANDARDS FOR CHECKOUT:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Transparent Language: Clearly display: 'By subscribing, you authorize Sri Sri Yoga to charge ₹[Amount] every [Cycle] starting on [Date] until cancelled.'\n• No Pre-Ticked Checkboxes: Consent boxes must be actively checked by user or explicitly confirmed via the CTA button.\n• Linked Policies: Clickable links to [Terms of Service], [Privacy Policy] and [Health Waiver] visible adjacent to CTA.\n• Order Summary Breakdown: Displays Initial Charge, Plan Duration, Renewal Amount, and Next Billing Date."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c24_r = add_card(s24, 6.8, 1.75, 5.7, 5.1, RGBColor(255, 251, 235), GOLD)
    tf = c24_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "SAMPLE VISUAL CONSENT COPY:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p = tf.add_paragraph()
    p.text = "┌──────────────────────────────────────────────┐\n│ ☑️ I agree to the Terms of Service, Privacy  │\n│    Policy, and Online Yoga Health Waiver.   │\n│                                              │\n│ 🔒 100% Secure Payment via Juspay Gateway    │\n│    Cancel or Pause anytime in 1-Click       │\n└──────────────────────────────────────────────┘"
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 25: Refund Policy Framework (Indian Wellness Standards)
    s25 = prs.slides.add_slide(blank_layout)
    set_bg(s25, prs)
    add_header(s25, "Consumer Protection", "Refund Policy Framework & Regulatory Guidelines (India)", "Balancing consumer fairness, digital service non-reversibility, and payment gateway SLAs")
    c25_l = add_card(s25, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c25_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "RECOMMENDED REFUND POLICY MATRIX:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Pre-Commencement Cancellation: 100% refund if cancelled before attending the first scheduled live session.\n• Duplicate / Accidental Debit: 100% automatic refund processed in 5-7 business days via Juspay.\n• Technical Non-Delivery: Full refund if live stream fails due to platform server downtime without makeup slots.\n• Post-Commencement: Non-refundable once live classes have been attended; user is offered Pause or Family Transfer option."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c25_r = add_card(s25, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c25_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "DISPUTE & CHARGEBACK MANAGEMENT:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "• 24-Hour Resolution SLA: Dedicated support queue for billing issues.\n• Juspay Gateway Webhook Logging: Comprehensive audit trail of payment timestamps, mandate tokens, and delivery notifications for dispute defense.\n• Proactive Goodwill Policy: Unused subscriptions can be converted into course credits for other Art of Living programs."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 26: Subscription State Machine Architecture
    s26 = prs.slides.add_slide(blank_layout)
    set_bg(s26, prs)
    add_header(s26, "System Architecture", "Complete Subscription State Machine & Lifecycle Transitions", "Formal lifecycle state diagram governing customer access, billing webhooks, and retry logic")
    states = [
        ("1. TRIAL", "14-day free access. ₹0 charge. Full class access."),
        ("2. ACTIVE", "Paid membership. Daily Zoom links & dashboard active."),
        ("3. PAUSED", "Classes frozen. End date extended. Booking triggers unpause."),
        ("4. GRACE_PERIOD", "Payment failed on renewal. 7-day retry grace. Access maintained."),
        ("5. CANCELLED", "Auto-renew stopped. Access continues until cycle end."),
        ("6. EXPIRED", "Subscription ended. Prompts 1-click repurchase.")
    ]
    for i, (st_name, st_desc) in enumerate(states):
        c = add_card(s26, 0.8 + (i%3)*3.95, 1.75 + (i//3)*2.6, 3.8, 2.35, BG_WHITE, CARD_BORDER)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = st_name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK
        p = tf.add_paragraph()
        p.text = st_desc
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(6)

    # Slide 27: Admin Console & Operations Capabilities
    s27 = prs.slides.add_slide(blank_layout)
    set_bg(s27, prs)
    add_header(s27, "Operations Console", "Admin Operations Console & Payment Traceability Ledger", "Comprehensive toolset for subscription management, renewals, retry triggers, and Juspay logs")
    c27_l = add_card(s27, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c27_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "ADMIN OPERATIONAL MODULES BUILT:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "1. Executive KPI Dashboard: Real-time active subscribers, monthly revenue, renewal success rates, and operational alerts.\n2. Subscriptions Management: Search by phone/name, filter by tier, view pause history, and trigger manual actions.\n3. Renewals & Failed Retries: Track dunning grace periods and trigger manual retry webhooks.\n4. Course & Plan Management: Configure batch slots, prices, and pause pools.\n5. Communications Center: Preview automated WhatsApp templates."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c27_r = add_card(s27, 6.8, 1.75, 5.7, 5.1, RGBColor(240, 249, 255), ACCENT_BLUE)
    tf = c27_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "PAYMENTS & MANDATES TRACEABILITY LEDGER:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p = tf.add_paragraph()
    p.text = "• End-to-End Audit Trail: Customer ➔ Subscription ➔ Upgrade Event ➔ Payment ➔ Mandate Token.\n• Advanced Filter Drawer: Filter by Course Plan, Payment Method (UPI / Card / NetBanking), Status (Paid / Failed), and Transaction Type (Initial / Renewal / Upgrade).\n• Real-Time Webhook Logs: Inspect Juspay webhook payloads and signature verification status."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 28: Subscription System Requirements Specification
    s28 = prs.slides.add_slide(blank_layout)
    set_bg(s28, prs)
    add_header(s28, "Product Requirements", "Subscription System Requirements Specification (PRD)", "Summary of functional requirements across User Experience, Admin Console, and Backend Systems")
    req_cols = [
        ("User Capabilities", "• Mobile/WhatsApp OTP login\n• 3-tap pre-pay questionnaire\n• 1-click Juspay checkout\n• Post-pay progressive profile\n• Self-serve pause & resume\n• 1-click plan upgrade\n• Transparent cancellation", PRIMARY),
        ("Admin Capabilities", "• Executive operational KPIs\n• Subscriber search & drawers\n• Payment ledger traceability\n• Failed renewal retry triggers\n• Manual pause overrides\n• Refund processing hooks\n• Course batch scheduler", ACCENT_BLUE),
        ("System Capabilities", "• Juspay webhook handlers\n• RBI pre-debit notifications\n• Dynamic date extension math\n• Dunning retry cadence\n• WhatsApp reminder bot\n• DPDP Act data encryption\n• SSO entitlement sync", SECONDARY)
    ]
    for i, (req_t, req_b, col) in enumerate(req_cols):
        c = add_card(s28, 0.8 + i*3.95, 1.75, 3.8, 5.1, BG_WHITE, CARD_BORDER)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
        p = tf.paragraphs[0]
        p.text = req_t
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col
        p = tf.add_paragraph()
        p.text = req_b
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # Slide 29: Phased Rollout Roadmap (Phase 1 to Phase 5)
    s29 = prs.slides.add_slide(blank_layout)
    set_bg(s29, prs)
    add_header(s29, "Implementation Roadmap", "5-Phase Rollout Roadmap & Implementation Timeline", "Structured deployment plan from MVP launch to advanced AI habit coaching")
    phases = [
        ("Phase 1: MVP Core (W1-4)", "• 'Pay First -> Profile After' flow\n• 3-question survey\n• Juspay 1-Time payment\n• Basic post-pay profile setup"),
        ("Phase 2: Lifecycle (W5-8)", "• Pause / Resume entitlement (15/30/45d)\n• Self-serve cancellation\n• Admin payments ledger & drawers"),
        ("Phase 3: Hybrid (W9-12)", "• Optional AutoPay mandate toggle\n• RBI pre-debit notifications\n• Upgrades & proration credit"),
        ("Phase 4: Optimization (W13-16)", "• 50/50 A/B checkout split testing\n• WhatsApp attendance bot\n• Dunning recovery automation"),
        ("Phase 5: Scale (W17+)", "• Multi-course companion bundles\n• Ashram loyalty rewards\n• AI personalized posture tips")
    ]
    for i, (p_t, p_desc) in enumerate(phases):
        c = add_card(s29, 0.8 + (i%3)*3.95, 1.75 + (i//3)*2.6, 3.8, 2.35, BG_WHITE, CARD_ORANGE_BORDER if i==0 else CARD_BORDER)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = p_t
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK if i==0 else TEXT_DARK
        p = tf.add_paragraph()
        p.text = p_desc
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(6)

    # Slide 30: A/B Testing Plan (Control vs Variant)
    s30 = prs.slides.add_slide(blank_layout)
    set_bg(s30, prs)
    add_header(s30, "Testing Protocol", "A/B Testing Protocol: Validating Conversion Uplift", "Controlled 50/50 split test to prove the conversion superiority of the 'Pay First' architecture")
    c30_l = add_card(s30, 0.8, 1.75, 5.7, 5.1, ALERT_RED_BG, RGBColor(254, 202, 202))
    tf = c30_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "CONTROL (50% Traffic) - BASELINE FLOW:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "• Landing -> Plan -> OTP -> 4 Questions -> 8-Field Registration Form -> Payment\n• Hypothesis: High drop-off at registration form prior to payment."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c30_r = add_card(s30, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c30_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "VARIANT (50% Traffic) - PAY FIRST FLOW:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "• Landing -> Plan -> OTP -> 3 Questions -> PAYMENT -> Post-Pay Profile\n• Expected Uplift: +25% to +40% increase in completed paid transactions.\n• Target Significance: 95% statistical confidence over 4-week test."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # Slide 31: Executive Recommendation: The 10 Key Decisions
    s31 = prs.slides.add_slide(blank_layout)
    set_bg(s31, prs)
    add_header(s31, "Executive Decisions", "Executive Recommendations: Answering the 10 Core Strategic Questions", "Clear, actionable leadership answers backed by UX research and competitive benchmarks")
    c31_l = add_card(s31, 0.8, 1.75, 5.7, 5.1, BG_WHITE, CARD_BORDER)
    tf = c31_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "CORE PRODUCT & PAYMENT DECISIONS:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "1. Move to Pay-First? YES. Supreme design rule to maximize conversion.\n2. Pre-Payment Info? Mobile OTP + 3 quick one-tap survey questions only.\n3. Post-Payment Info? Full Name, Email, Age, City, Batch, and Health notes.\n4. 1-Time vs Recurring? Default to 1-Time Upfront; offer AutoPay as optional toggle.\n5. ~20% Renewal Strategy? Implement habit metrics and automated WhatsApp cadence."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)

    c31_r = add_card(s31, 6.8, 1.75, 5.7, 5.1, RGBColor(255, 247, 237), CARD_ORANGE_BORDER)
    tf = c31_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "LIFECYCLE & OPERATIONAL DECISIONS:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "6. Offer Pause/Resume? YES. 15d (3M), 30d (6M), 45d (12M) with early resume refunds.\n7. Cancellation UX? Transparent 1-click self-serve cancellation (anti-dark patterns).\n8. Pricing Architecture? ₹1,499 (3M), ₹2,699 (6M), ₹4,999 (12M Annual Best Value).\n9. 45-Day Access? Cloud entitlement token with day countdown and upgrade bridge.\n10. MVP Launch Plan? Phase 1 Pay-First MVP launch within 4 weeks."
    p.font.size = Pt(9.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)

    # Slide 32: Conclusion & Next Steps
    s32 = prs.slides.add_slide(blank_layout)
    set_bg(s32, prs, RGBColor(255, 248, 232))
    c32 = add_card(s32, 0.8, 1.1, 11.733, 5.3, BG_WHITE, CARD_ORANGE_BORDER)
    tf = c32.text_frame
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.6)
    p = tf.paragraphs[0]
    p.text = "THE ART OF LIVING • SRI SRI YOGA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "Ready for Stakeholder Approval & Production Build"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)
    p = tf.add_paragraph()
    p.text = "• Working prototype is live and fully interactive on local staging server.\n• All customer flows (Payment First, Survey, Pause, Unpause, Manage Hub) tested and verified.\n• Ready to proceed with Juspay production merchant keys and CRM API integration."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(12)

    # Save Deck 1
    master_strategy_filename = "Sri_Sri_Yoga_Subscription_UX_Master_Deck.pptx"
    prs.save(master_strategy_filename)
    print(f"Successfully generated 32-slide Strategy Deck: {master_strategy_filename}")

    # =========================================================================
    # DECK 2: 10-SLIDE DEDICATED LEGAL & COMPLIANCE PRESENTATION
    # =========================================================================
    prs_leg = pptx.Presentation()
    prs_leg.slide_width = Inches(13.333)
    prs_leg.slide_height = Inches(7.5)

    def add_leg_header(slide, category, title, subtitle=None):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY_DARK

        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.55))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        if subtitle:
            s_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.28), Inches(11.7), Inches(0.35))
            tf_s = s_box.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
            p_s = tf_s.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(11.5)
            p_s.font.color.rgb = TEXT_MUTED

    # Legal Slide 1: Cover
    ls1 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls1, prs_leg, RGBColor(248, 250, 252))
    lc1 = add_card(ls1, 0.8, 1.1, 11.733, 5.3, BG_WHITE, CARD_BORDER)
    tf = lc1.text_frame
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.6)
    p = tf.paragraphs[0]
    p.text = "THE ART OF LIVING • SRI SRI YOGA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "Legal, Regulatory & Compliance Framework"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)
    p = tf.add_paragraph()
    p.text = "RBI e-Mandate Guidelines, Explicit Consent, DPDP Act 2023, Refund Policies & Consumer Safeguards"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "Prepared for: Legal Counsel & Risk Compliance Teams | September 2026"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p.space_before = Pt(32)

    # Legal Slide 2: RBI e-Mandate Framework
    ls2 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls2, prs_leg)
    add_leg_header(ls2, "RBI Regulations", "RBI e-Mandate Framework & Pre-Debit Notification Rules", "Statutory compliance requirements for recurring payments under Reserve Bank of India circulars")
    c_l2 = add_card(ls2, 0.8, 1.75, 11.733, 5.1)
    tf = c_l2.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "MANDATORY RBI COMPLIANCE PROVISIONS (AFA & e-MANDATE):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "1. Additional Factor of Authentication (AFA): The initial recurring mandate registration must be authenticated via AFA (OTP / UPI PIN).\n\n2. Mandatory Pre-Debit Notifications: A pre-debit notification via SMS/Email must be delivered to the customer at least 24 to 48 hours before the actual auto-renewal debit occurs.\n\n3. Opt-Out & Unlink Facility: Customers must have an explicit option to pause or revoke the mandate without requiring customer support intervention.\n\n4. Transaction Limits: Transactions up to ₹15,000 per cycle do not require recurring OTP if valid e-mandate exists.\n\n5. Juspay Role: Juspay handles tokenization, webhook dispatches and pre-debit notification triggering."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 3: Explicit User Consent Standards
    ls3 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls3, prs_leg)
    add_leg_header(ls3, "User Consent", "Explicit User Consent & Payment Authorization Architecture", "Ensuring clear, unambiguous consent without deceptive pre-selected checkboxes")
    c_l3 = add_card(ls3, 0.8, 1.75, 11.733, 5.1)
    tf = c_l3.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "EXPLICIT CONSENT STANDARDS FOR RECURRING BILLING:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Transparent Language: Clearly display: 'By subscribing, you authorize Sri Sri Yoga to charge ₹[Amount] every [Cycle] starting on [Date] until cancelled.'\n• No Pre-Ticked Checkboxes: Consent boxes must be actively checked by user or explicitly confirmed via the CTA button.\n• Clear Distinction: If recurring payment is optional, the toggle between '1-Time Purchase' and 'Auto-Renew & Save' must be distinct.\n• Linked Policies: Clickable links to [Terms of Service] and [Privacy Policy] must be visible adjacent to the checkout button.\n• Order Summary Display: Breakdown showing Initial Charge, Trial Duration, Renewal Amount, and Next Billing Date."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 4: Terms of Service & Online Yoga Health Waiver
    ls4 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls4, prs_leg)
    add_leg_header(ls4, "Terms & Privacy", "Terms of Service & Online Yoga Health Waiver Framework", "Key clauses governing online yoga instruction, health waivers, and personal data protection (DPDP Act)")
    c_l4 = add_card(ls4, 0.8, 1.75, 11.733, 5.1)
    tf = c_l4.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "KEY CLAUSES FOR ONLINE YOGA SUBSCRIBER TERMS:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "1. Health & Physical Activity Waiver: Explicit acknowledgment that participant is in suitable physical health for asanas and assumes responsibility for personal injury.\n\n2. Intellectual Property & Recording: Prohibition against recording, rebroadcasting or redistributing live stream sessions.\n\n3. Dynamic Schedule Alteration: Art of Living reserves right to substitute certified teachers or adjust session schedules with prior notice.\n\n4. Digital Personal Data Protection (DPDP Act 2023): Personal data (phone, email, health notes) stored securely with AES-256 encryption and processed solely for yoga delivery."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 5: Statutory Cancellation & Mandate Revocation Rights
    ls5 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls5, prs_leg)
    add_leg_header(ls5, "Cancellation Rights", "Statutory Cancellation & Mandate Revocation Rights", "Ensuring frictionless self-serve cancellation in full compliance with Consumer Protection Rules")
    c_l5 = add_card(ls5, 0.8, 1.75, 11.733, 5.1)
    tf = c_l5.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "MANDATORY CANCELLATION & REVOCATION PROVISIONS:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Self-Serve Accessibility: Customers must be able to cancel auto-renewal in <= 3 clicks within the web portal without human agent intervention.\n• Revocation Webhook: Cancellation triggers instant mandate revocation API call to Juspay and NPCI.\n• Continuation of Paid Validity: Paid access continues undisturbed until the end of the paid period.\n• Written Confirmation: Instant dispatch of cancellation confirmation with effective date via Email and WhatsApp."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 6: Refund Policy Matrix
    ls6 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls6, prs_leg)
    add_leg_header(ls6, "Refund Policy", "Comprehensive Refund Policy Framework (Indian Wellness)", "Balancing consumer fairness, digital service non-reversibility, and payment gateway SLAs")
    c_l6 = add_card(ls6, 0.8, 1.75, 11.733, 5.1)
    tf = c_l6.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "RECOMMENDED REFUND POLICY MATRIX (SUBJECT TO LEGAL SIGN-OFF):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Pre-Commencement Cancellation: 100% refund if requested before the first scheduled class date.\n• Duplicate / Accidental Debit: 100% refund processed within 5-7 business days via Juspay gateway.\n• Technical Non-Delivery: Full refund if live class streaming fails due to server outages without alternative slots.\n• Post-Commencement Subscription: Non-refundable after active class attendance, but users can PAUSE their membership or transfer validity to another family member.\n• Chargeback Handling: Automated webhook logging provides full audit trail for disputed transactions."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 7: Free Trial Disclosures
    ls7 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls7, prs_leg)
    add_leg_header(ls7, "Trial Disclosures", "Free Trial Disclosures & Auto-Debit Transition Safeguards", "Ensuring transparent communication during 14-day trial to auto-renew transition")
    c_l7 = add_card(ls7, 0.8, 1.75, 11.733, 5.1)
    tf = c_l7.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "TRIAL TRANSITION COMPLIANCE PROTOCOL:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Day 0 Signup Disclosure: Clearly states '14 Days Free, then ₹4,999/year starting on [Date]. Cancel anytime during trial for ₹0 charge.'\n• Day 11 Advance Alert (WhatsApp & Email): 'Your free trial ends in 3 days. Your annual subscription will renew on [Date].'\n• Day 13 Final 24h Reminder: 'Your Sri Sri Yoga annual membership will activate tomorrow. Manage or cancel here: [Link]'\n• Day 14 Mandate Execution: Official GST tax invoice dispatched immediately upon successful charge."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 8: Data Privacy & DPDP Act 2023
    ls8 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls8, prs_leg)
    add_leg_header(ls8, "Data Protection", "Digital Personal Data Protection (DPDP Act 2023) Compliance", "Data minimization, purpose limitation, and subscriber consent architecture")
    c_l8 = add_card(ls8, 0.8, 1.75, 11.733, 5.1)
    tf = c_l8.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "DPDP ACT 2023 COMPLIANCE MANDATES:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "1. Data Minimization: Pre-payment questionnaire captures only high-level wellness intent. Sensitive personal data collected post-payment only for yoga delivery.\n2. Purpose Limitation: Subscriber data used exclusively for course coordination and ashram updates.\n3. Right to Erasure: Subscribers can request deletion of account and health records via privacy portal.\n4. Consent Artifact Logging: Consent timestamp and IP address stored in secure audit log."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 9: Legal Sign-off Checklist
    ls9 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls9, prs_leg)
    add_leg_header(ls9, "Sign-Off Checklist", "Legal & Operational Sign-Off Checklist", "Required validation items before production deployment of Juspay AutoPay integration")
    c_l9 = add_card(ls9, 0.8, 1.75, 11.733, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c_l9.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "LEGAL SIGN-OFF CHECKLIST PRIOR TO LIVE LAUNCH:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "☑️ Terms of Service & Online Yoga Health Waiver reviewed by legal counsel\n☑️ Privacy Policy updated with DPDP Act 2023 compliance\n☑️ RBI e-Mandate pre-debit SMS/Email notification template approved\n☑️ Refund & Cancellation Policy published on website footer and checkout modal\n☑️ Juspay AutoPay mandate registration terms verified against NPCI guidelines\n☑️ Customer Support Escalation SLA (24h turnaround) established for billing disputes"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 10: Legal Contacts & Governance
    ls10 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls10, prs_leg, RGBColor(248, 250, 252))
    lc10 = add_card(ls10, 0.8, 1.1, 11.733, 5.3, BG_WHITE, CARD_BORDER)
    tf = lc10.text_frame
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.6)
    p = tf.paragraphs[0]
    p.text = "THE ART OF LIVING • SRI SRI YOGA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "Legal & Compliance Approval Protocol"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)
    p = tf.add_paragraph()
    p.text = "• Legal Deck ready for review by Art of Living Legal Counsel & Finance Director.\n• All compliance artifacts, terms templates, and consent copies are finalized for production sign-off."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(12)

    master_legal_filename = "Sri_Sri_Yoga_Legal_Compliance_Master_Deck.pptx"
    prs_leg.save(master_legal_filename)
    print(f"Successfully generated 10-slide Legal & Compliance Deck: {master_legal_filename}")

if __name__ == '__main__':
    create_master_presentation_decks()
