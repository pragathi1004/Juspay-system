import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    PRIMARY = RGBColor(234, 88, 12)       # Saffron / Orange (#ea580c)
    PRIMARY_LIGHT = RGBColor(255, 237, 213)# Light Saffron (#ffedd5)
    PRIMARY_DARK = RGBColor(154, 52, 18)  # Deep Saffron (#9a3412)
    SECONDARY = RGBColor(22, 101, 52)     # Forest Green (#166534)
    SECONDARY_LIGHT = RGBColor(220, 252, 231) # Light Green (#dcfce7)
    ACCENT_TEAL = RGBColor(13, 148, 136)  # Teal (#0d9488)
    TEXT_DARK = RGBColor(30, 41, 59)      # Slate 800 (#1e293b)
    TEXT_MUTED = RGBColor(100, 116, 139)  # Slate 500 (#64748b)
    BG_CREAM = RGBColor(253, 251, 247)    # Warm Cream (#fdfbf7)
    BG_WHITE = RGBColor(255, 255, 255)    # Pure White
    CARD_BORDER = RGBColor(226, 232, 240) # Slate 200
    CARD_ORANGE_BORDER = RGBColor(254, 215, 170)
    ALERT_RED = RGBColor(220, 38, 38)     # Red (#dc2626)
    ALERT_RED_BG = RGBColor(254, 242, 242)# Light Red (#fef2f2)

    def set_bg(slide, color=BG_CREAM):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=None):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY

        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.55))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(21)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        if subtitle:
            s_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.35))
            tf_s = s_box.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
            p_s = tf_s.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(12)
            p_s.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, bg_color=BG_WHITE, border_color=CARD_BORDER, radius_adjust=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, RGBColor(255, 248, 232))

    top_band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = PRIMARY
    top_band.line.fill.background()

    c1 = add_card(s1, 0.8, 1.1, 11.733, 5.3, BG_WHITE, CARD_ORANGE_BORDER)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.8)
    tf1.margin_top = Inches(0.6)
    tf1.margin_right = Inches(0.8)

    p = tf1.paragraphs[0]
    p.text = "THE ART OF LIVING • SRI SRI YOGA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    p = tf1.add_paragraph()
    p.text = "Subscription & Auto-Payment UX Strategy"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    p = tf1.add_paragraph()
    p.text = "Research-Backed Recommendations for Purchase Funnel Optimization, Payment Models & Subscription Lifecycle Management"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(8)

    p = tf1.add_paragraph()
    p.text = "CORE STRATEGIC MANDATES & DECISIONS:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p.space_before = Pt(28)

    p = tf1.add_paragraph()
    p.text = "• Core Principle: PAY FIRST → COMPLETE PROFILE AFTER (Eliminating the Pre-Payment Registration Trap)\n• Business Reality: Evaluating One-Time vs Recurring Payments in light of ~20% Current Renewal Rates\n• Competitor Benchmarks: Multi-dimensional study of Cult.fit, Habuild, Satvic Movement, Isha Yoga, Times Health+, Kamya\n• Complete Subscription Lifecycle: Pause entitlement, frictionless cancellation, plan switching & teacher credibility"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 2: EXECUTIVE SUMMARY & STRATEGIC PILLARS
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "Executive Decision Mandate", "Executive Summary: 3 Transformational Decisions", "Key strategic choices to maximize purchase conversion and long-term spiritual community retention")

    cards_data_s2 = [
        ("1. FUNNEL FRICTION", PRIMARY, "Pay First → Profile After", 
         "• Current blocker: 8+ mandatory profile fields (Age, PIN, City) demanded BEFORE payment.\n• Psychological impact: Creates severe friction at peak purchase intent.\n• Strategic Action: Shift personal profile collection post-payment. Collect ONLY authentication + payment essentials pre-checkout.\n• Expected Impact: +25% to +40% increase in checkout completions."),
        ("2. PAYMENT MODEL", ACCENT_TEAL, "Evaluate ~20% Renewal Reality",
         "• Current context: Baseline renewal rate is ~20%.\n• Risk of forced auto-debit: High mandate failure rates, customer anxiety, cancellation friction.\n• Strategic Action: Position 1-Time payment as default / transparent choice, with recurring auto-debit offered with explicit perks/discounts.\n• Expected Impact: Maximize initial conversion without compromising user trust."),
        ("3. RETENTION & TRUST", SECONDARY, "Lifecycle UX & Teacher Trust",
         "• Pause Entitlement: Adopt 15/30/45-day pause pools to prevent cancellations during travel/illness.\n• Clean Cancellation: Zero dark patterns; transparent self-serve cancellation.\n• Teacher Lineage: Showcase teacher profiles & credentials on landing page to build spiritual credibility.")
    ]

    for idx, (pill, color, title, body) in enumerate(cards_data_s2):
        c = add_card(s2, 0.8 + idx * 3.98, 1.75, 3.77, 5.1)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
        p = tf.paragraphs[0]
        p.text = pill
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(6)
        p = tf.add_paragraph()
        p.text = body
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 3: CORE BUSINESS PROBLEM (FRICTION PEAK)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "Problem Diagnosis", "The Core Business Problem: The Pre-Payment Registration Trap", "Why asking for extensive profile details before payment suppresses checkout conversions")

    c3_left = add_card(s3, 0.8, 1.75, 5.7, 5.1, ALERT_RED_BG, RGBColor(252, 165, 165))
    tf = c3_left.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "⚠️ THE CURRENT FRICTION PEAK"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "Information Overload Before Value Realization"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "In the current prototype, immediately after WhatsApp OTP login, users are presented with a heavy 8-field form demanding:\n\n• First Name & Last Name\n• WhatsApp Number (repeat)\n• Email Address\n• Age & Date of Birth\n• Postal / PIN Code\n• City & State\n• Class Language Selection\n• Multiple checkbox consents\n\nConsequence: The user has demonstrated purchase intent, but is forced into administrative data entry before experiencing any product value or completing checkout."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    c3_right = add_card(s3, 6.8, 1.75, 5.7, 5.1)
    tf = c3_right.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "📉 CONVERSION DROP-OFF CHAIN"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "How Cognitive Load Kills Payment Conversion"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)

    steps = [
        ("1. High Effort Required", "Typing PIN, age, name on mobile keyboard causes mental resistance."),
        ("2. Privacy & Hesitation", "Users question why Sri Sri Yoga needs PIN code/age before they even pay."),
        ("3. Fatigue & Drop-Off", "Users abandon the tab/browser before ever reaching the payment gateway."),
        ("4. Lost Revenue", "Profile completion is treated as more important than payment capture.")
    ]
    for stitle, sdesc in steps:
        p = tf.add_paragraph()
        p.text = f"• {stitle}: {sdesc}"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # Save progress
    prs.save("Daily_Yoga_Subscription_UX_Strategy_Updated.pptx")
    print("Built through Slide 3.")

build_presentation()
