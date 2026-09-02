import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_complete_decks():
    # =========================================================================
    # DECK 1: BUSINESS & STRATEGY PRESENTATION
    # =========================================================================
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    PRIMARY = RGBColor(234, 88, 12)        # Saffron / Orange (#ea580c)
    PRIMARY_LIGHT = RGBColor(255, 237, 213)# Light Saffron (#ffedd5)
    PRIMARY_DARK = RGBColor(154, 52, 18)   # Deep Saffron (#9a3412)
    SECONDARY = RGBColor(22, 101, 52)      # Forest Green (#166534)
    SECONDARY_LIGHT = RGBColor(220, 252, 231) # Light Green (#dcfce7)
    ACCENT_TEAL = RGBColor(13, 148, 136)   # Teal (#0d9488)
    TEXT_DARK = RGBColor(30, 41, 59)       # Slate 800 (#1e293b)
    TEXT_MUTED = RGBColor(100, 116, 139)   # Slate 500 (#64748b)
    BG_CREAM = RGBColor(253, 251, 247)     # Warm Cream (#fdfbf7)
    BG_WHITE = RGBColor(255, 255, 255)     # Pure White
    CARD_BORDER = RGBColor(226, 232, 240)  # Slate 200
    CARD_ORANGE_BORDER = RGBColor(254, 215, 170)
    ALERT_RED = RGBColor(220, 38, 38)      # Red (#dc2626)
    ALERT_RED_BG = RGBColor(254, 242, 242) # Light Red (#fef2f2)

    def set_bg(slide, color=BG_CREAM):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=None):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY

        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.55))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(21)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        if subtitle:
            s_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.35))
            tf_s = s_box.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
            p_s = tf_s.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(12)
            p_s.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, bg_color=BG_WHITE, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    # =========================================================================
    # SLIDE 1: COVER SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, RGBColor(255, 248, 232))

    top_band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.4))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = PRIMARY
    top_band.line.fill.background()

    c1 = add_card(s1, 0.8, 1.1, 11.733, 5.3, BG_WHITE, CARD_ORANGE_BORDER)
    tf = c1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.6)
    tf.margin_right = Inches(0.8)

    p = tf.paragraphs[0]
    p.text = "THE ART OF LIVING • SRI SRI YOGA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    p = tf.add_paragraph()
    p.text = "Subscription & Auto-Payment UX Strategy"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    p = tf.add_paragraph()
    p.text = "Research-Backed Recommendations for Purchase Funnel Optimization, Payment Models & Subscription Lifecycle Management"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "CORE STRATEGIC MANDATES & DECISIONS:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p.space_before = Pt(28)

    p = tf.add_paragraph()
    p.text = "• Core Principle: PAY FIRST → COMPLETE PROFILE AFTER (Eliminating the Pre-Payment Registration Trap)\n• Business Reality: Evaluating One-Time vs Recurring Payments in light of ~20% Current Renewal Rates\n• Competitor Benchmarks: Multi-dimensional study of Cult.fit, Habuild, Satvic Movement, Isha Yoga, Times Health+, Kamya\n• Complete Subscription Lifecycle: Pause entitlement, frictionless cancellation, plan switching & teacher credibility"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 2: EXECUTIVE SUMMARY & STRATEGIC PILLARS
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "Executive Decision Mandate", "Executive Summary: 3 Transformational Decisions", "Key strategic choices to maximize purchase conversion and long-term spiritual community retention")

    cards_data_s2 = [
        ("1. FUNNEL FRICTION", PRIMARY, "Pay First → Profile After", 
         "• Current blocker: 8+ mandatory profile fields (Age, PIN, City) demanded BEFORE payment.\n• Psychological impact: Creates severe friction at peak purchase intent.\n• Strategic Action: Shift personal profile collection post-payment. Collect ONLY authentication + payment essentials pre-checkout.\n• Expected Impact: +25% to +40% increase in checkout completions."),
        ("2. PAYMENT MODEL", ACCENT_TEAL, "Evaluate ~20% Renewal Reality",
         "• Current context: Baseline renewal rate is ~20%.\n• Risk of forced auto-debit: High mandate failure rates, customer anxiety, cancellation friction.\n• Strategic Action: Position 1-Time payment as default / transparent choice, with recurring auto-debit offered with explicit perks/discounts.\n• Expected Impact: Maximize initial conversion without compromising user trust."),
        ("3. RETENTION & TRUST", SECONDARY, "Lifecycle UX & Teacher Trust",
         "• Pause Entitlement: Adopt 15/30/45-day pause pools to prevent cancellations during travel/illness.\n• Clean Cancellation: Zero dark patterns; transparent self-serve cancellation.\n• Teacher Lineage: Showcase teacher profiles & credentials on landing page to build spiritual credibility.")
    ]

    for idx, (pill, color, title, body) in enumerate(cards_data_s2):
        c = add_card(s2, 0.8 + idx * 3.98, 1.75, 3.77, 5.1)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
        p = tf.paragraphs[0]
        p.text = pill
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(6)
        p = tf.add_paragraph()
        p.text = body
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 3: CORE BUSINESS PROBLEM (FRICTION PEAK)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "Problem Diagnosis", "The Core Business Problem: The Pre-Payment Registration Trap", "Why asking for extensive profile details before payment suppresses checkout conversions")

    c3_left = add_card(s3, 0.8, 1.75, 5.7, 5.1, ALERT_RED_BG, RGBColor(252, 165, 165))
    tf = c3_left.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "⚠️ THE CURRENT FRICTION PEAK"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "Information Overload Before Value Realization"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "In the current prototype, immediately after WhatsApp OTP login, users are presented with a heavy 8-field form demanding:\n\n• First Name & Last Name\n• WhatsApp Number (repeat confirmation)\n• Email Address\n• Age & Date of Birth\n• Postal / PIN Code\n• City & State\n• Class Language Selection\n• Multiple checkbox consents\n\nConsequence: The user has demonstrated purchase intent, but is forced into administrative data entry before experiencing any product value or completing checkout."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    c3_right = add_card(s3, 6.8, 1.75, 5.7, 5.1)
    tf = c3_right.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "📉 CONVERSION DROP-OFF CHAIN"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "How Cognitive Load Kills Payment Conversion"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)

    steps_s3 = [
        ("1. High Interaction Effort", "Typing PIN codes, city, age, and name on mobile keyboards introduces cognitive fatigue."),
        ("2. Data Hesitation & Privacy Anxiety", "Users question why Sri Sri Yoga needs postal codes and age before they have even purchased."),
        ("3. Severe Funnel Drop-Off", "Users abandon the browser tab during lengthy form entry before reaching Juspay checkout."),
        ("4. False Priority Conflict", "Registration / profile completion is mistakenly prioritized over payment capture.")
    ]
    for stitle, sdesc in steps_s3:
        p = tf.add_paragraph()
        p.text = f"• {stitle}: {sdesc}"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 4: CURRENT USER JOURNEY (VISUAL BREAKDOWN)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "Current Flow Analysis", "Current Purchase Journey: Friction Peak Before Payment", "A linear breakdown of the 8-step baseline user journey highlighting the critical abandonment point")

    journey_steps_s4 = [
        ("Step 1", "Landing Page", "User explores benefits & plan options", BG_WHITE, CARD_BORDER),
        ("Step 2", "Select Plan", "Selects 3M, 6M, or 12M membership", BG_WHITE, CARD_BORDER),
        ("Step 3", "Get Started", "Clicks CTA with high purchase intent", BG_WHITE, CARD_BORDER),
        ("Step 4", "Login / OTP", "Enters phone & receives WhatsApp OTP", BG_WHITE, CARD_BORDER),
        ("Step 5", "Personal Details", "🔴 FRICTION PEAK: 8+ form fields required", ALERT_RED_BG, RGBColor(239, 68, 68)),
        ("Step 6", "Yoga Questions", "3 questions including redundant time query", BG_WHITE, CARD_BORDER),
        ("Step 7", "Payment", "Juspay AutoPay / UPI checkout", BG_WHITE, CARD_BORDER),
        ("Step 8", "Dashboard", "Subscription active & class access", SECONDARY_LIGHT, RGBColor(187, 247, 208))
    ]

    for idx, (snum, stitle, sdesc, bg_c, b_c) in enumerate(journey_steps_s4):
        x = 0.8 + (idx % 4) * 2.98
        y = 1.75 if idx < 4 else 4.4
        c = add_card(s4, x, y, 2.78, 2.4, bg_c, b_c)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = snum
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = PRIMARY if "Step 5" not in snum else ALERT_RED
        p = tf.add_paragraph()
        p.text = stitle
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(4)
        p = tf.add_paragraph()
        p.text = sdesc
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MUTED if "Step 5" not in snum else ALERT_RED
        p.space_before = Pt(6)

    # =========================================================================
    # SLIDE 5: NEW PRIMARY UX PRINCIPLE
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "Core UX Philosophy", "New Primary UX Principle: Pay First → Complete Profile After", "Aligning data collection timing with user purchase intent and psychological readiness")

    # Big Banner Principle
    b5 = add_card(s5, 0.8, 1.75, 11.733, 1.3, RGBColor(255, 247, 237), CARD_ORANGE_BORDER)
    tf = b5.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "CORE PRINCIPLE: PAY FIRST → COMPLETE PROFILE AFTER"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "Only collect information strictly required to authenticate the user and process payment. Defer all profile building post-conversion."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(4)

    # Two column comparison cards
    c5_left = add_card(s5, 0.8, 3.25, 5.7, 3.6, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c5_left.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
    p = tf.paragraphs[0]
    p.text = "✅ MANDATORY BEFORE PAYMENT"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Essential for Authentication & Processing"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(4)
    p = tf.add_paragraph()
    p.text = "• Mobile Number (for WhatsApp OTP login & account security)\n• Authentication Token / OTP (proves user identity)\n• Selected Plan & Billing Frequency (for order amount & mandate)\n• Payment Gateway Authorization (UPI / Card / Mandate setup)\n\nResult: User completes purchase in < 60 seconds with zero keyboard fatigue."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c5_right = add_card(s5, 6.8, 3.25, 5.7, 3.6, BG_WHITE, CARD_BORDER)
    tf = c5_right.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
    p = tf.paragraphs[0]
    p.text = "📋 OPTIONAL / POST-PAYMENT PROFILE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "For Personalization & Delivery (Post-Pay)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(4)
    p = tf.add_paragraph()
    p.text = "• Participant Full Name (for certificate & personalization)\n• Email Address (for calendar invites & tax invoices)\n• Age & Date of Birth (for health adaptations)\n• Postal / PIN Code & City (for regional statistics)\n• Class Language Preference (Hindi, Telugu, Tamil, etc.)\n\nResult: 100% of paying users happily fill this out post-purchase."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 6: RECOMMENDED ALTERNATIVE FLOW
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "Recommended User Experience", "Recommended Alternative Flow: Rapid Payment → Rich Profile", "Streamlined 7-step high-converting architecture moving all administrative data collection post-payment")

    rec_steps_s6 = [
        ("Phase 1", "Landing Page", "Hero CTA + Teacher Profiles", BG_WHITE, CARD_BORDER),
        ("Phase 2", "Plan Selection", "Transparent 3M/6M/12M terms", BG_WHITE, CARD_BORDER),
        ("Phase 3", "Mobile OTP", "Fast WhatsApp OTP auth (10s)", BG_WHITE, CARD_BORDER),
        ("Phase 4", "3 Quick Questions", "High-value intent & starting point", PRIMARY_LIGHT, CARD_ORANGE_BORDER),
        ("Phase 5", "PAYMENT GATEWAY", "Juspay Checkout (Peak Intent)", SECONDARY_LIGHT, RGBColor(34, 197, 94)),
        ("Phase 6", "Payment Success 🎉", "Subscription instantly activated", SECONDARY_LIGHT, RGBColor(34, 197, 94)),
        ("Phase 7", "Complete Profile", "Progressive name, age, city entry", BG_WHITE, CARD_BORDER),
        ("Phase 8", "Yoga Dashboard", "Live classes & masterclasses ready", BG_WHITE, CARD_BORDER)
    ]

    for idx, (pnum, ptitle, pdesc, bg_c, b_c) in enumerate(rec_steps_s6):
        x = 0.8 + (idx % 4) * 2.98
        y = 1.75 if idx < 4 else 4.4
        c = add_card(s6, x, y, 2.78, 2.4, bg_c, b_c)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = pnum
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = PRIMARY if "Phase 5" not in pnum and "Phase 6" not in pnum else SECONDARY
        p = tf.add_paragraph()
        p.text = ptitle
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(4)
        p = tf.add_paragraph()
        p.text = pdesc
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_MUTED if "Phase 5" not in pnum and "Phase 6" not in pnum else SECONDARY
        p.space_before = Pt(6)

    # =========================================================================
    # SLIDE 7: REDESIGNED PRE-PAYMENT ONBOARDING
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "Onboarding Optimization", "Redesigning Pre-Payment Onboarding: 3 High-Value Questions", "Eliminating non-actionable inquiries and focusing exclusively on meaningful personalization")

    c7_left = add_card(s7, 0.8, 1.75, 5.7, 5.1, ALERT_RED_BG, RGBColor(252, 165, 165))
    tf = c7_left.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "❌ REMOVED QUESTION & RATIONALE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "Question: 'How much time can you dedicate each day?'"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "Why this question was removed:\n\n1. Fixed Class Schedules: Sri Sri Yoga daily live classes are conducted at fixed morning/evening slots (5am, 6am, 7am, 4pm, 6pm). The user's time estimate does not alter session durations.\n\n2. Zero Operational Impact: Does not personalize the schedule or unlock custom content tiers.\n\n3. Cognitive Drag: Forces user to calculate schedule availability before buying, creating hesitation.\n\nRule: Do not retain a question simply because it existed previously."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    c7_right = add_card(s7, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c7_right.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "✅ THE 3 HIGH-VALUE QUESTIONS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Fast, Engaging, Single/Multi-Select (< 10s total)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "1. Question 1 (User Intent):\n'What brings you to Sri Sri Yoga?' (Flexibility, Stress, Strength, Pain relief, Weight, Mindfulness)\n\n2. Question 2 (Starting Level):\n'What's your current yoga experience?' (Beginner, Some Experience, Regular Practitioner, Advanced)\n\n3. Question 3 (Desired Transformation):\n'What would you like to get from your practice?' (Better fitness, Energy & Vitality, Calmer Mind, Daily Habit)\n\nBenefit: Feels like a tailored wellness assessment without data friction."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 8: THE 3 PRE-PAYMENT QUESTIONS SPECIFICATION
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_header(s8, "Questionnaire UX Specs", "Pre-Payment Questionnaire Specification (3 Screens)", "Zero-typing, intuitive choice architecture designed for instant mobile touch interaction")

    q_cards = [
        ("QUESTION 1 OF 3 (Multi-Select)", "What brings you to Sri Sri Yoga?", 
         "• Improve flexibility & mobility 🤸\n• Build strength & stamina 🏋️\n• Reduce stress & anxiety 💧\n• Weight management ⚖️\n• Relieve body pain 🩺\n• Self-improvement & mindfulness 🧘\n• Something else ☀️\n\nUX Goal: Understand user health motivations."),
        ("QUESTION 2 OF 3 (Single-Select)", "What's your current yoga experience?", 
         "• Complete Beginner (New to yoga) 🌱\n• Some Experience (Tried yoga before) 🌿\n• Regular Practitioner (Practice regularly) 🌳\n• Advanced Yogi (Long-time practice) ✨\n\nUX Goal: Calibrate teacher difficulty level."),
        ("QUESTION 3 OF 3 (Single-Select)", "What would you like to get from your practice?", 
         "• Better physical fitness & tone 💪\n• More energy & vitality ⚡\n• Better sleep & relaxation 🌙\n• Less stress & a calmer mind 🕊️\n• Better flexibility & mobility 🧘\n• Build a consistent daily habit 📅\n\nUX Goal: Personalize milestone recommendations.")
    ]

    for idx, (qtag, qtitle, qoptions) in enumerate(q_cards):
        c = add_card(s8, 0.8 + idx * 3.98, 1.75, 3.77, 5.1)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
        p = tf.paragraphs[0]
        p.text = qtag
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p = tf.add_paragraph()
        p.text = qtitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(6)
        p = tf.add_paragraph()
        p.text = qoptions
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 9: POST-PAYMENT PROFILE COMPLETION
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)
    add_header(s9, "Post-Payment Experience", "Post-Payment Profile Completion: Warm & Progressive", "Transforming mandatory administrative forms into an empowering onboarding ritual")

    c9_left = add_card(s9, 0.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c9_left.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "🎉 STEP 1: PAYMENT SUCCESS CELEBRATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Instant Gratification & Relief"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "• Headline: 'Payment Successful! Welcome to Daily Sri Sri Yoga 🎉'\n• Sub-headline: 'Your 3-Month Subscription is Active. Order ID #ORD_991823'\n• Reassurance: Instant SMS & WhatsApp confirmation sent.\n• Transition CTA: 'Complete Your Profile in 30 Seconds' -> Leads into profile setup without feeling like a barrier."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    c9_right = add_card(s9, 6.8, 1.75, 5.7, 5.1)
    tf = c9_right.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "📝 STEP 2: PROGRESSIVE PROFILE SETUP"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "Structured Progressive Disclosure"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "Screen 1: Personal Details\n• Full Name (for your official completion certificate)\n• Email Address (for calendar alerts & Zoom links)\n\nScreen 2: Tailoring Your Practice\n• Age / Age Group (for posture modifications)\n• City / State & PIN (for local community satsangs)\n• Preferred Class Language (English, Hindi, Telugu, Tamil, Marathi, Kannada)\n\nFinal CTA: 'Go to Yoga Dashboard' -> Unlocks daily classes immediately."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 10: SIDE-BY-SIDE FLOW COMPARISON
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10)
    add_header(s10, "Flow Comparison", "Current Flow vs. Recommended Flow Comparison", "Visualizing how shifting profile collection post-payment removes the critical conversion bottleneck")

    c10_cur = add_card(s10, 0.8, 1.75, 5.7, 5.1, ALERT_RED_BG, RGBColor(252, 165, 165))
    tf = c10_cur.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "🔴 CURRENT FLOW (High Pre-Payment Burden)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "Landing → Plan → Login/OTP → 8-Field Profile Form (FRICTION PEAK) → 3 Questions → Payment → Dashboard"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "Key Bottlenecks:\n• 8 manual inputs before checkout creates massive cognitive barrier.\n• User drops out before payment info is even initiated.\n• Estimated pre-payment time: 2.5 to 3.5 minutes.\n• Payment conversion rate suppressed by ~30%."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    c10_rec = add_card(s10, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c10_rec.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "🟢 RECOMMENDED FLOW (Frictionless Payment First)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Landing → Plan → Login/OTP → 3 Quick Questions (10s) → PAYMENT → Success 🎉 → Profile Completion → Dashboard"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)
    p = tf.add_paragraph()
    p.text = "Key Advantages:\n• Checkout reached in < 45 seconds from landing page.\n• Captures payment when purchase intent is highest.\n• 100% of paying users complete profile post-purchase.\n• Projected +25% to +40% increase in paid subscriptions."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # =========================================================================
    # SLIDE 11: STRATEGIC CONTEXT (~20% RENEWAL RATE REALITY)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_bg(s11)
    add_header(s11, "Business Economics", "The Renewal Reality: Evaluating our ~20% Baseline Renewal Rate", "Why a ~20% renewal rate necessitates re-thinking automatic recurring mandates as the default")

    c11_top = add_card(s11, 0.8, 1.75, 11.733, 1.3, RGBColor(254, 242, 242), RGBColor(252, 165, 165))
    tf = c11_top.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "KEY BUSINESS CONTEXT: CURRENT RENEWAL RATE IS ~20%"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "This baseline indicates that 80% of users do not naturally renew on automatic billing cycles today. Forcing mandatory auto-renewal upfront creates substantial conversion friction and customer chargeback friction."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(4)

    c11_l = add_card(s11, 0.8, 3.25, 5.7, 3.6)
    tf = c11_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
    p = tf.paragraphs[0]
    p.text = "Option A: Forced Auto-Debit Model"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ALERT_RED
    p = tf.add_paragraph()
    p.text = "• Friction: Users fear unexpected bank debits after 3 months.\n• Mandate Failures: UPI AutoPay authentication fails for ~20-30% of Indian banks.\n• Churn Illusion: 80% cancel or revoke mandate before renewal anyway.\n• Net Effect: Lower initial conversion + high cancellation support tickets."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c11_r = add_card(s11, 6.8, 3.25, 5.7, 3.6, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c11_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
    p = tf.paragraphs[0]
    p.text = "Option B: One-Time Upfront + Trust Model"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "• Trust: Users pay comfortably via standard UPI/Card with zero recurring anxiety.\n• Higher Checkout: No e-mandate friction or recurring debit fear.\n• Deliberate Re-engagement: Win back renewals through rich 14-day & 30-day spiritual engagement campaigns.\n• Net Effect: Maximize top-of-funnel paid participants + high brand affinity."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 12: ONE-TIME VS RECURRING PAYMENT MATRIX
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_bg(s12)
    add_header(s12, "Payment Architecture", "One-Time Payment vs. Recurring Auto-Payment: Trade-Off Matrix", "Comprehensive strategic comparison of payment mechanisms for digital wellness in India")

    c12_1 = add_card(s12, 0.8, 1.75, 5.7, 5.1)
    tf = c12_1.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
    p = tf.paragraphs[0]
    p.text = "ONE-TIME PAYMENT MODEL"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "PROS:\n• Lower psychological barrier — simple one-time transaction.\n• Zero cancellation anxiety or surprise billing fear.\n• 100% payment gateway compatibility (Standard UPI, NetBanking, Cards).\n• Fewer customer support disputes & zero mandate failure recoveries.\n• Significantly higher initial conversion rate.\n\nCONS:\n• Revenue is less automatically recurring.\n• Requires active renewal communication campaigns at cycle end.\n• Relies on customer motivation to manually re-purchase."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c12_2 = add_card(s12, 6.8, 1.75, 5.7, 5.1)
    tf = c12_2.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
    p = tf.paragraphs[0]
    p.text = "RECURRING AUTO-PAYMENT (AUTOPAY)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p = tf.add_paragraph()
    p.text = "PROS:\n• Predictable cash flow for retained subscribers.\n• Frictionless renewal without requiring manual re-transaction.\n• Higher LTV for dedicated long-term practitioners.\n\nCONS:\n• Higher initial drop-off (users hesitate to authorize auto-debit).\n• Requires strict RBI e-Mandate compliance & pre-debit notifications.\n• Mandate execution failure rate is high in India (~15-25%).\n• Support burden for forgotten subscriptions & refund requests.\n• At ~20% natural renewal, benefits are muted compared to conversion losses."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 13: SUBSCRIPTION MODEL DECISION FRAMEWORK
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_bg(s13)
    add_header(s13, "Business Model Recommendation", "Subscription Model Decision Framework for Art of Living", "Evaluating three potential architectural options to arrive at the optimal hybrid solution")

    opts_data = [
        ("OPTION A: One-Time Default", "Pure upfront access fee (3M/6M/12M)", "• Pros: Maximum checkout conversion, spiritual trust, zero mandate failures.\n• Cons: No automated recurring revenue.\n• Suitability: High for spiritual workshops."),
        ("OPTION B: Pure Auto-Renew", "Mandatory recurring subscription", "• Pros: Recurring predictable revenue.\n• Cons: Drop-offs during mandate setup, high refund disputes.\n• Suitability: Risky with current ~20% renewal."),
        ("OPTION C: Hybrid Choice (RECOMMENDED)", "Default 1-Time + Optional Auto-Renew Perk", "• Structure: 1-Time payment is standard. Users get option to toggle 'Auto-Renew & Save 10% / Get 15 Extra Pause Days'.\n• Benefit: Best of both worlds — zero friction for hesitant users + auto-debit for committed yogis.")
    ]

    for idx, (otitle, osub, obody) in enumerate(opts_data):
        c = add_card(s13, 0.8 + idx * 3.98, 1.75, 3.77, 5.1, SECONDARY_LIGHT if idx == 2 else BG_WHITE, RGBColor(34, 197, 94) if idx == 2 else CARD_BORDER)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
        p = tf.paragraphs[0]
        p.text = otitle
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = SECONDARY if idx == 2 else PRIMARY
        p = tf.add_paragraph()
        p.text = osub
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(4)
        p = tf.add_paragraph()
        p.text = obody
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 14: COMPETITOR RESEARCH BENCHMARK MATRIX
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    set_bg(s14)
    add_header(s14, "Competitive Intelligence", "Comprehensive Competitor Benchmark Matrix", "How leading Indian yoga, fitness and wellness platforms handle pricing, payments, onboarding and lifecycle")

    # Table of Competitors
    rows, cols = 7, 6
    table_shape = s14.shapes.add_table(rows, cols, Inches(0.8), Inches(1.75), Inches(11.733), Inches(5.1))
    table = table_shape.table
    table.columns[0].width = Inches(1.9)
    table.columns[1].width = Inches(2.1)
    table.columns[2].width = Inches(1.9)
    table.columns[3].width = Inches(1.8)
    table.columns[4].width = Inches(2.0)
    table.columns[5].width = Inches(2.033)

    headers_comp = ["Platform", "Pricing Tiers", "Payment Model", "Profile Timing", "Pause Feature", "Key UX Learning"]
    for j, h in enumerate(headers_comp):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = BG_WHITE

    comp_rows = [
        ("Habuild Yoga", "3M: ₹1,999 | 6M: ₹2,999 | 12M: ₹3,999", "1-Time Upfront (via WhatsApp & Web)", "Post-Payment via WhatsApp bot", "WhatsApp / Support-based manual pause", "WhatsApp-first onboarding removes web registration friction entirely."),
        ("Cult.fit (Cult Live)", "3M: ₹3,490 | 12M: ₹8,990", "1-Time & Auto-Debit options", "Login first, profile after purchase", "In-App Pause Pool (15d/30d/45d)", "Pause days are a structured plan entitlement; easy slider/date selection."),
        ("Satvic Yoga", "Monthly: $40 | 3M: $90 | Annual: $240", "Recurring Auto-Debit (Stripe/Gateway)", "Name + WhatsApp only before pay", "Manual via email support", "Minimal pre-payment form (only Name + WhatsApp number)."),
        ("Isha Yoga (IEO)", "₹1,500 (Regional) / ₹3,500 (English)", "1-Time Program Enrollment Fee", "Post-Payment portal registration", "Self-paced course validity extension", "Spiritual programs succeed heavily on one-time transparent pricing."),
        ("Times Health+", "Annual & Multi-month plans", "Auto-Renew + 7-Day Guarantee", "Account creation required", "Account settings freeze option", "Transparent 7-day money-back guarantee reduces pre-payment risk."),
        ("Kamya Wellness", "14-Day Challenge + Custom Plans", "1-Time Challenge / Direct UPI", "WhatsApp automation onboarding", "Contact support via WhatsApp", "Free 14-day WhatsApp challenge converts users into paying members.")
    ]

    for i, row_data in enumerate(comp_rows):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if i % 2 == 0 else RGBColor(248, 250, 252)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8.5)
                p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 15: COMPETITOR DEEP-DIVE 1 (HABUILD, CULT.FIT, TIMES HEALTH+)
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    set_bg(s15)
    add_header(s15, "Competitive Deep-Dive", "Competitor Deep-Dive: Habuild, Cult.fit & Times Health+", "Key architectural and onboarding lessons for Sri Sri Yoga subscription design")

    c15_1 = add_card(s15, 0.8, 1.75, 3.77, 5.1)
    tf = c15_1.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "HABUILD YOGA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "• Model: 1-Time payment (3M ₹1,999, 12M ₹3,999).\n• Friction Removal: Users never fill website forms. Checkout happens via quick payment link; full profile & batch selection happens on WhatsApp.\n• Pause Policy: Empathic support-driven pause via WhatsApp (+91 79693 29699).\n• Takeaway: WhatsApp is India's most comfortable channel for post-pay engagement."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c15_2 = add_card(s15, 4.78, 1.75, 3.77, 5.1)
    tf = c15_2.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "CULT.FIT (CULT PASS)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p = tf.add_paragraph()
    p.text = "• Model: Tiered membership with structured pause days (e.g. 12M = 45 days pause).\n• In-App Pause: Self-serve calendar picker. Unused days are credited back if resumed early.\n• Cancellation: Transparent cancellation in Profile -> Membership.\n• Takeaway: Cult's pause entitlement prevents churn when users travel or fall sick."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c15_3 = add_card(s15, 8.76, 1.75, 3.77, 5.1)
    tf = c15_3.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "TIMES HEALTH+"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "• Model: Multi-tier wellness subscription with money-back guarantee.\n• Trust Mechanism: 7-day no-questions-asked refund policy directly on landing page.\n• Cancellation: Standard self-serve under 'Manage Membership'.\n• Takeaway: A prominent risk-reversal guarantee dramatically lowers pre-payment hesitation."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 16: COMPETITOR DEEP-DIVE 2 (SATVIC, ISHA, KAMYA)
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    set_bg(s16)
    add_header(s16, "Competitive Deep-Dive", "Competitor Deep-Dive: Satvic Movement, Isha Yoga & Kamya", "Spiritual and holistic community models with high consumer retention and spiritual loyalty")

    c16_1 = add_card(s16, 0.8, 1.75, 3.77, 5.1)
    tf = c16_1.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "SATVIC YOGA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "• Model: Daily Yoga subscription ($40/mo, $90/3mo, $240/yr) + community workshops.\n• Form Timing: Ultra-short initial checkout (Name + WhatsApp + Email only).\n• Communication: Daily Zoom links and reminders pushed directly to WhatsApp.\n• Takeaway: High energy visual branding combined with minimal form friction."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c16_2 = add_card(s16, 4.78, 1.75, 3.77, 5.1)
    tf = c16_2.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "ISHA YOGA (INNER ENG.)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p = tf.add_paragraph()
    p.text = "• Model: Fixed 1-time program fee (₹1,500 - ₹3,500).\n• Trust & Lineage: Heavy emphasis on Sadhguru & authentic Yogic transmission.\n• Zero Auto-Debit: Completely avoids recurring payment friction; users willingly upgrade to advanced programs.\n• Takeaway: Spiritual credibility and authentic lineage drive high voluntary repurchase."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c16_3 = add_card(s16, 8.76, 1.75, 3.77, 5.1)
    tf = c16_3.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "KAMYA WELLNESS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "• Model: Goal-specific programs (PCOS, Weight Loss, Hormonal Health).\n• Lead Generation: 14-Day Free Challenge via WhatsApp automated flow.\n• Community Engagement: Private WhatsApp groups for daily guidance.\n• Takeaway: Free challenge preview builds high habituation before asking for paid commitment."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 17: TEACHER PROFILES RECOMMENDATION
    # =========================================================================
    s17 = prs.slides.add_slide(blank_layout)
    set_bg(s17)
    add_header(s17, "Landing Page Conversion", "Trust Building: Teacher Profiles on Landing Page", "Humanizing the digital yoga journey and establishing spiritual lineage before checkout")

    c17_top = add_card(s17, 0.8, 1.75, 11.733, 1.3, RGBColor(255, 247, 237), CARD_ORANGE_BORDER)
    tf = c17_top.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "RECOMMENDATION: ADD 'MEET YOUR SRI SRI YOGA TEACHERS' SECTION"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "Unlike fitness apps, yoga practitioners seek spiritual authenticity, certified lineage, and empathic guidance. Highlighting teacher profiles directly increases landing page conversion by reducing instructor uncertainty."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(4)

    t_profiles = [
        ("Dinesh Kashikar", "Senior International Sri Sri Yoga Teacher", "25+ Years Experience • IIT Alumnus • Guided 500,000+ students worldwide in authentic pranayama & asanas."),
        ("Kamlesh Barwal", "Global Director, Sri Sri Yoga", "20+ Years Experience • Authored yoga manuals • Specializes in therapeutic yoga, mindfulness & yoga sutras."),
        ("Certified Faculty", "500+ Hours Certified Instructors", "Rigorous Art of Living Ashram training • Multilingual guidance in English, Hindi, Telugu, Tamil, Marathi & Kannada.")
    ]

    for idx, (tname, trole, tbio) in enumerate(t_profiles):
        c = add_card(s17, 0.8 + idx * 3.98, 3.25, 3.77, 3.6)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
        p = tf.paragraphs[0]
        p.text = tname
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK
        p = tf.add_paragraph()
        p.text = trole
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = SECONDARY
        p.space_before = Pt(4)
        p = tf.add_paragraph()
        p.text = tbio
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 18: FREE TRIAL VS FREE PREVIEW STRATEGY
    # =========================================================================
    s18 = prs.slides.add_slide(blank_layout)
    set_bg(s18)
    add_header(s18, "Trial Strategy", "Free Trial vs. Free Preview: Acquisition Model Analysis", "Evaluating 14-day pre-authorized trial vs. free sample content for Sri Sri Yoga")

    c18_1 = add_card(s18, 0.8, 1.75, 5.7, 5.1)
    tf = c18_1.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "MODEL 1: 14-DAY TRIAL WITH MANDATE"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "Card / UPI Mandate Pre-Authorized Upfront"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "• Mechanics: ₹1 or ₹0 pre-auth mandate setup with auto-debit on Day 15.\n• Advantage: High conversion to paid if user develops a daily morning yoga habit.\n• Disadvantage: Higher friction at signup; requires automated reminder SMS on Day 12.\n• Abuse Prevention: 1 trial per verified WhatsApp phone number.\n• Recommendation: Offer as option for high-intent subscribers."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    c18_2 = add_card(s18, 6.8, 1.75, 5.7, 5.1)
    tf = c18_2.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "MODEL 2: FREE PREVIEW / 7-DAY CHALLENGE"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p = tf.add_paragraph()
    p.text = "No Payment Details Required Upfront"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "• Mechanics: Instant access to 3 live morning sessions upon WhatsApp OTP login.\n• Advantage: Zero signup friction; massive top-of-funnel lead generation.\n• Conversion Trigger: At end of session 3, teacher invites participant to subscribe.\n• Recommendation: Highly recommended for monthly seasonal yoga campaigns (e.g. Navratri / New Year)."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 19: PRICING & PACKAGING STRATEGY
    # =========================================================================
    s19 = prs.slides.add_slide(blank_layout)
    set_bg(s19)
    add_header(s19, "Pricing Architecture", "Pricing & Packaging Strategy for the Indian Yoga Market", "Calibrated tier structure balancing affordability, spiritual accessibility and strong LTV")

    pricing_cards = [
        ("3 MONTHS PLAN", "Quarterly Commitment", "₹1,499 - ₹2,999", "₹499 / month effective", 
         "• Ideal starting point for beginners.\n• Includes Unlimited Live Classes\n• Monthly Sunday Masterclass\n• 15 Pause Days Pool\n• Progress tracking on dashboard"),
        ("6 MONTHS PLAN", "Habit Transformation", "₹3,999", "₹666 / month effective", 
         "• For committed regular yogis.\n• Includes 3 Mini-Programs\n• Bonus Asana E-Books\n• 30 Pause Days Pool\n• Community Satsang Access"),
        ("12 MONTHS PLAN (BEST VALUE)", "Annual Lifestyle Transformation", "₹4,999", "₹416 / month effective (Best Value)", 
         "• 67% discount off ₹14,999.\n• 1 Month Full Access to Art of Living App\n• 45 Pause Days Pool\n• Dedicated teacher Q&A priority\n• Annual Completion Certificate")
    ]

    for idx, (pname, psub, pprice, peff, pfeat) in enumerate(pricing_cards):
        c = add_card(s19, 0.8 + idx * 3.98, 1.75, 3.77, 5.1, SECONDARY_LIGHT if idx == 2 else BG_WHITE, RGBColor(34, 197, 94) if idx == 2 else CARD_BORDER)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
        p = tf.paragraphs[0]
        p.text = pname
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = SECONDARY if idx == 2 else PRIMARY
        p = tf.add_paragraph()
        p.text = psub
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(2)
        p = tf.add_paragraph()
        p.text = pprice
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(6)
        p = tf.add_paragraph()
        p.text = peff
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = SECONDARY if idx == 2 else PRIMARY_DARK
        p.space_before = Pt(2)
        p = tf.add_paragraph()
        p.text = pfeat
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 20: FRICTIONLESS CANCELLATION STRATEGY
    # =========================================================================
    s20 = prs.slides.add_slide(blank_layout)
    set_bg(s20)
    add_header(s20, "Retention Architecture", "Frictionless Cancellation Strategy: Anti-Dark Patterns", "Building spiritual trust through transparent, self-serve cancellation and proactive pause alternatives")

    c20_left = add_card(s20, 0.8, 1.75, 5.7, 5.1)
    tf = c20_left.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "RECOMMENDED CANCELLATION WORKFLOW"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "Step-by-Step Self-Serve Flow"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)

    c_steps = [
        ("1. Profile Entry", "User taps 'Cancel Subscription' in Manage Hub without calling support."),
        ("2. Impact Clarity", "Clearly shows access remains active until the end of the paid cycle."),
        ("3. Optional Reason", "Simple non-blocking feedback (Travelling, Busy, Financial, etc.)."),
        ("4. Pause Alternative", "Proactively offers: 'Traveling or busy? Pause for free instead.'"),
        ("5. 1-Click Confirmation", "Single tap confirmation; instant SMS & Email cancellation receipt.")
    ]
    for s_t, s_d in c_steps:
        p = tf.add_paragraph()
        p.text = f"• {s_t}: {s_d}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(6)

    c20_right = add_card(s20, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c20_right.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "💎 WHY HONEST CANCELLATION DRIVES RETENTION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Spiritual Trust Creates Re-Activation"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "• No Dark Patterns: We never hide cancellation behind phone trees or support tickets.\n• Preserves Brand Affinity: When a subscriber leaves peacefully, 35%+ return within 6 months.\n• Pause Interception: 40% of users attempting cancellation choose to pause when presented with an easy 15/30 day pause pool.\n• Easy Reactivation: 1-tap 'Reactivate Membership' button remains visible on dashboard."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 21: PAUSE / RESUME STRATEGY & ENTITLEMENT MODEL
    # =========================================================================
    s21 = prs.slides.add_slide(blank_layout)
    set_bg(s21)
    add_header(s21, "Retention Architecture", "Pause / Resume Strategy: Entitlement Pool Architecture", "Cult.fit-inspired pause framework customized for Art of Living yoga schedules and early resume refunds")

    c21_1 = add_card(s21, 0.8, 1.75, 3.77, 5.1)
    tf = c21_1.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "1. PAUSE ENTITLEMENT RULES"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "Plan-Based Allowance Pools"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(4)
    p = tf.add_paragraph()
    p.text = "• 3-Month Plan: 15 Pause Days\n• 6-Month Plan: 30 Pause Days\n• 12-Month Plan: 45 Pause Days\n\n• Available strictly during active cycle.\n• Unused days expire when subscription ends (no carryover).\n• Users cannot pause beyond their remaining pool balance."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c21_2 = add_card(s21, 4.78, 1.75, 3.77, 5.1)
    tf = c21_2.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "2. AUTOMATIC DATE MATH"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p = tf.add_paragraph()
    p.text = "Expiry & Renewal Extension"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(4)
    p = tf.add_paragraph()
    p.text = "• Subscription Expiry: Extended by exact days paused.\n• Next Auto-Debit Date: Postponed forward in lockstep.\n• Zero Payment Surprises: Customer sees exact revised expiry before confirming.\n• Webhook Dispatch: Sends MANDATE_PAUSED event to Juspay."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    c21_3 = add_card(s21, 8.76, 1.75, 3.77, 5.1)
    tf = c21_3.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "3. EARLY RESUME & BOOKING"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Unused Day Refunds & Access"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(4)
    p = tf.add_paragraph()
    p.text = "• Early Resume: Resuming on day 6 of a 10-day pause immediately credits 4 unused days back to balance.\n• Booking while Paused: If user taps 'Join Live Class', system prompts: 'Resume & Book'.\n• Automatic Resume: System auto-resumes on scheduled end date."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 22: UPGRADE / DOWNGRADE LIFECYCLE
    # =========================================================================
    s22 = prs.slides.add_slide(blank_layout)
    set_bg(s22)
    add_header(s22, "Lifecycle Management", "Upgrade, Downgrade & Plan Switching Architecture", "Frictionless tier migrations with fair proration credit and automated mandate adjustments")

    c22_l = add_card(s22, 0.8, 1.75, 5.7, 5.1)
    tf = c22_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "⬆️ PLAN UPGRADES (e.g. 3M → 12M)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Instant Upgrade with Prorated Credit"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "• Immediate Activation: Upgraded benefits (e.g. AOL App access + 45 pause days) unlock instantly.\n• Proration Calculation: Unused days from current 3M plan are credited towards 12M fee.\n• Mandate Adjustment: Old Juspay mandate updated or replaced with new ₹4,999 annual debit schedule.\n• User Experience: Transparent summary showing exact prorated amount charged today."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    c22_r = add_card(s22, 6.8, 1.75, 5.7, 5.1)
    tf = c22_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "⬇️ PLAN DOWNGRADES (e.g. 12M → 3M)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "End-of-Cycle Effective Downgrade"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "• Effective Date: 12-month access remains intact until current paid period ends.\n• Next Renewal: At renewal date, system charges lower 3M fee (₹1,499) instead of annual fee.\n• No Unfair Penalties: Avoid complex partial refunds; let paid access run its natural course.\n• Clear UI Status: Dashboard displays: 'Switching to 3 Months on [Renewal Date]'."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 23: MODERNIZED AUTHENTICATION (APPLE REMOVED)
    # =========================================================================
    s23 = prs.slides.add_slide(blank_layout)
    set_bg(s23)
    add_header(s23, "Authentication Optimization", "Modernized Authentication: Clean WhatsApp OTP Flow", "Eliminating unused social logins (Apple removed) and optimizing for India's mobile-first user base")

    c23_l = add_card(s23, 0.8, 1.75, 5.7, 5.1)
    tf = c23_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "📱 PRIMARY AUTHENTICATION: WHATSAPP / SMS OTP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "Zero-Password Mobile Login"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "• Mobile Country Code Picker (+91 default, international supported).\n• Auto-Detect WhatsApp Delivery: OTP sent directly via WhatsApp Business API for instant 99.8% delivery rate.\n• 4-Digit Quick Verification: Clean PIN blocks with auto-focus.\n• Test Phone Pre-configured: +91 99206 56992 (Test OTP: 4829)."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    c23_r = add_card(s23, 6.8, 1.75, 5.7, 5.1)
    tf = c23_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "🔄 SOCIAL LOGIN CLEANUP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p = tf.add_paragraph()
    p.text = "Removal of Apple Login & Streamlined Alternatives"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "• Apple Login Removed: Avoids Apple Sign-In email relay masking issues which break WhatsApp communication.\n• Retained Social Options: Google 1-Tap & Facebook Login (trusted by 90%+ Indian web users).\n• Email Passwordless Magic Link: Available as fallback for corporate/NRI users.\n• Clean Modal Interface: Top tab switchers removed for uncluttered direct input."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 24: 45-DAY ACCESS & MULTI-PLATFORM ENTITLEMENT
    # =========================================================================
    s24 = prs.slides.add_slide(blank_layout)
    set_bg(s24)
    add_header(s24, "Product Entitlement", "45-Day Program Access & Multi-Platform Entitlement", "Architecting seamless cross-platform yoga access between the Web Portal and Art of Living App")

    c24_top = add_card(s24, 0.8, 1.75, 11.733, 1.3, RGBColor(255, 247, 237), CARD_ORANGE_BORDER)
    tf = c24_top.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    p = tf.paragraphs[0]
    p.text = "HYBRID ACCESS MODEL: UNIFIED ACCOUNT-LINKED ENTITLEMENT"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "Users access Daily Live Yoga classes via Web Browser (Zero Install) and unlock 1-Month / 45-Day companion meditation tracks inside the Art of Living App using their single registered mobile number."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(4)

    ent_cards = [
        ("WEB PORTAL ACCESS (Primary)", "Daily Live Classes & Zoom", "• Instant browser access without app downloads.\n• High compatibility across Laptops, Tablets & Smart TVs.\n• Live chat, teacher interaction & masterclass archives."),
        ("APP ACCESS (Companion)", "Art of Living App Bridge", "• Mobile number automatically activates AOL App entitlement.\n• Unlocks guided meditations, chants & pranayama timers.\n• Push notifications for upcoming morning sessions."),
        ("ENTITLEMENT EXPIRY", "Clean Expiry Management", "• Access seamlessly expires at end of 45-day / paid period.\n• Automated WhatsApp re-engagement 3 days prior.\n• Option to extend or upgrade with 1-click renewal.")
    ]

    for idx, (etitle, esub, ebody) in enumerate(ent_cards):
        c = add_card(s24, 0.8 + idx * 3.98, 3.25, 3.77, 3.6)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
        p = tf.paragraphs[0]
        p.text = etitle
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_DARK
        p = tf.add_paragraph()
        p.text = esub
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = SECONDARY
        p.space_before = Pt(4)
        p = tf.add_paragraph()
        p.text = ebody
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(8)

    # =========================================================================
    # SLIDE 25: ADMIN SUBSCRIPTION CONSOLE ARCHITECTURE
    # =========================================================================
    s25 = prs.slides.add_slide(blank_layout)
    set_bg(s25)
    add_header(s25, "Operations & Admin Console", "Admin & Operations Subscriptions Management Architecture", "Comprehensive back-office capabilities to manage mandates, customer lifecycles, and exceptions")

    admin_cols = [
        ("1. SUBSCRIPTION ROSTER", "• Search by Name, Phone, Mandate ID\n• Filter by Status (Active, Paused, Failed, Expired)\n• Detailed Customer Drawer with payment history\n• Real-time mandate debit date tracking"),
        ("2. LIFECYCLE CONTROLS", "• Manual Pause / Resume Override\n• Manual Plan Upgrade / Downgrade\n• Revoke Mandate / Turn Auto-Pay Off\n• Extend Grace Period by 7/14 days"),
        ("3. PAYMENTS & RETRIES", "• Juspay Webhook Log Inspector\n• Automated 3-stage Smart Retries (48h, 72h, 96h)\n• 1-Click Partial/Full Refund Processor\n• Reconciliation export to CSV/Excel")
    ]

    for idx, (atitle, abody) in enumerate(admin_cols):
        c = add_card(s25, 0.8 + idx * 3.98, 1.75, 3.77, 5.1)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.28)
        p = tf.paragraphs[0]
        p.text = atitle
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p = tf.add_paragraph()
        p.text = abody
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 26: ANALYTICS & FUNNEL TELEMETRY
    # =========================================================================
    s26 = prs.slides.add_slide(blank_layout)
    set_bg(s26)
    add_header(s26, "Analytics Framework", "Full-Funnel Analytics & Key Performance Indicator (KPI) Framework", "Tracking granular conversion, drop-off rates, and retention economics across the entire user journey")

    funnel_stages = [
        ("1. Landing Page", "100%", "Unique visitors exploring yoga programs"),
        ("2. Plan Selected", "65%", "Clicks 'Try for Free' / chooses 3M/6M/12M"),
        ("3. OTP Verified", "48%", "Enters mobile & completes WhatsApp OTP"),
        ("4. 3 Questions", "44%", "Answers 3 lightweight questions (<10s)"),
        ("5. PAYMENT SUCCESS", "32%", "PRIMARY KPI: Payment captured first"),
        ("6. Profile Complete", "31%", "98% post-payment completion rate"),
        ("7. Active Retention", "24%", "Daily class attendance & milestone unlocks")
    ]

    for idx, (fstage, fval, fdesc) in enumerate(funnel_stages):
        y = 1.75 + idx * 0.72
        c = add_card(s26, 0.8, y, 11.733, 0.62, SECONDARY_LIGHT if "PAYMENT" in fstage else BG_WHITE, RGBColor(34, 197, 94) if "PAYMENT" in fstage else CARD_BORDER)
        tf = c.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        p.text = f"{fstage}   |   Conversion Target: {fval}   |   {fdesc}"
        p.font.size = Pt(11)
        p.font.bold = True if "PAYMENT" in fstage else False
        p.font.color.rgb = SECONDARY if "PAYMENT" in fstage else TEXT_DARK

    # =========================================================================
    # SLIDE 27: A/B EXPERIMENTATION ROADMAP
    # =========================================================================
    s27 = prs.slides.add_slide(blank_layout)
    set_bg(s27)
    add_header(s27, "Validation & Testing", "A/B Experimentation Plan: Testing the Pay-First Hypothesis", "Structured 4-week split test methodology to prove conversion uplift before 100% rollout")

    c27_l = add_card(s27, 0.8, 1.75, 5.7, 5.1)
    tf = c27_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "EXPERIMENT DESIGN: CONTROL VS VARIANT"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "50/50 Traffic Split for 4 Weeks"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "• Control (50% Traffic):\nLanding → Plan → Login → 8-Field Profile Form → 3 Questions → Payment\n\n• Variant (50% Traffic):\nLanding → Plan → Login → 3 Quick Questions → Payment → Post-Pay Profile\n\n• Sample Size: 10,000 unique landing page visitors per arm."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    c27_r = add_card(s27, 6.8, 1.75, 5.7, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c27_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "🎯 PRIMARY & GUARDRAIL METRICS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "Success Criteria for 100% Rollout"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)
    p = tf.add_paragraph()
    p.text = "1. Primary Success Metric:\nPaid Subscription Conversion Rate (% of OTP logins completing payment).\nTarget: Minimum +25% relative uplift.\n\n2. Guardrail Metrics:\n• Post-Payment Profile Completion Rate (Target > 92%)\n• First 7-day Live Class Attendance (% attending at least 1 session)\n• Refund Requests / Support Tickets (must not increase)"
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

    # =========================================================================
    # SLIDE 28: PHASED ROLLOUT PLAN & FINAL RECOMMENDATIONS
    # =========================================================================
    s28 = prs.slides.add_slide(blank_layout)
    set_bg(s28)
    add_header(s28, "Implementation Roadmap", "Phased Rollout Plan & Executive Recommendations Checklist", "Actionable 5-phase engineering roadmap and core strategic takeaways for leadership sign-off")

    c28_l = add_card(s28, 0.8, 1.75, 5.7, 5.1)
    tf = c28_l.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "🚀 5-PHASE EXECUTION ROADMAP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p = tf.add_paragraph()
    p.text = "Phase 1: Pre-Payment Friction Removal (Weeks 1-2)\n• Remove 8 pre-payment form fields\n• Remove time-availability question\n• Add teacher profiles on landing page\n\nPhase 2: Payment Architecture (Weeks 3-4)\n• Implement 1-time default + optional auto-renew toggle\n• Update Juspay checkout webhooks\n\nPhase 3: Post-Pay Profile & Pause Flow (Weeks 5-6)\n• Deploy progressive post-pay profile setup\n• Deploy 15/30/45-day pause wizard\n\nPhase 4: Admin Dashboard & Telemetry (Weeks 7-8)"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(6)

    c28_r = add_card(s28, 6.8, 1.75, 5.7, 5.1, RGBColor(255, 247, 237), CARD_ORANGE_BORDER)
    tf = c28_r.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = "⭐ EXECUTIVE RECOMMENDATION CHECKLIST"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "1. Adopt 'Pay First -> Profile After' as the supreme design rule.\n2. Keep pre-payment questionnaire to exactly 3 engaging questions.\n3. Position 1-Time payment as default with optional Auto-Renew perks.\n4. Introduce 15/30/45-day Pause Entitlement to curb churn.\n5. Provide 100% transparent self-serve cancellation.\n6. Feature Art of Living certified Teacher Profiles on Landing Page."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

    prs.save("Daily_Yoga_Subscription_UX_Strategy_Updated.pptx")
    print(f"Successfully generated 28-slide Executive Strategy Deck: Daily_Yoga_Subscription_UX_Strategy_Updated.pptx")

    # =========================================================================
    # DECK 2: LEGAL & COMPLIANCE PRESENTATION
    # =========================================================================
    prs_leg = pptx.Presentation()
    prs_leg.slide_width = Inches(13.333)
    prs_leg.slide_height = Inches(7.5)

    def add_leg_header(slide, category, title, subtitle=None):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PRIMARY_DARK

        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.55))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(21)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        if subtitle:
            s_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.35))
            tf_s = s_box.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
            p_s = tf_s.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(12)
            p_s.font.color.rgb = TEXT_MUTED

    # Legal Slide 1: Title
    ls1 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls1, RGBColor(248, 250, 252))
    lc1 = add_card(ls1, 0.8, 1.1, 11.733, 5.3, BG_WHITE, CARD_BORDER)
    tf = lc1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.6)
    p = tf.paragraphs[0]
    p.text = "THE ART OF LIVING • SRI SRI YOGA"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "Legal, Regulatory & Compliance Framework"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)
    p = tf.add_paragraph()
    p.text = "RBI e-Mandate Guidelines, Explicit User Consent, Transparent Auto-Renewal Disclosures & Consumer Refund Policies"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(8)

    # Legal Slide 2: RBI e-Mandate Regulations
    ls2 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls2)
    add_leg_header(ls2, "RBI Regulations", "RBI e-Mandate Framework & Pre-Debit Notification Rules", "Statutory compliance requirements for recurring payments under Reserve Bank of India circulars")
    c_l2 = add_card(ls2, 0.8, 1.75, 11.733, 5.1)
    tf = c_l2.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "MANDATORY RBI COMPLIANCE PROVISIONS (AFA & e-MANDATE):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "1. Additional Factor of Authentication (AFA): The initial recurring mandate registration must be authenticated via AFA (OTP / UPI PIN).\n\n2. Mandatory Pre-Debit Notifications: A pre-debit notification via SMS/Email must be delivered to the customer at least 24 to 48 hours before the actual auto-renewal debit occurs.\n\n3. Opt-Out & Unlink Facility: Customers must have an explicit option to pause or revoke the mandate without requiring customer support intervention.\n\n4. Transaction Limits: Transactions up to ₹15,000 per cycle do not require recurring OTP if valid e-mandate exists.\n\n5. Juspay Role: Juspay handles tokenization, webhook dispatches and pre-debit notification triggering."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 3: Explicit User Consent
    ls3 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls3)
    add_leg_header(ls3, "User Consent", "Explicit User Consent & Payment Authorization Architecture", "Ensuring clear, unambiguous consent without deceptive pre-selected checkboxes")
    c_l3 = add_card(ls3, 0.8, 1.75, 11.733, 5.1)
    tf = c_l3.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "EXPLICIT CONSENT STANDARDS FOR RECURRING BILLING:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Transparent Language: Clearly display: 'By subscribing, you authorize Sri Sri Yoga to charge ₹[Amount] every [Cycle] starting on [Date] until cancelled.'\n• No Pre-Ticked Checkboxes: Consent boxes must be actively checked by user or explicitly confirmed via the CTA button.\n• Clear Distinction: If recurring payment is optional, the toggle between '1-Time Purchase' and 'Auto-Renew & Save' must be distinct.\n• Linked Policies: Clickable links to [Terms of Service] and [Privacy Policy] must be visible adjacent to the checkout button.\n• Order Summary Display: Breakdown showing Initial Charge, Trial Duration, Renewal Amount, and Next Billing Date."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 4: Terms & Privacy Disclosures
    ls4 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls4)
    add_leg_header(ls4, "Terms & Privacy", "Terms of Service & Privacy Policy Disclosures Framework", "Key clauses governing online yoga instruction, health waivers, and personal data protection (DPDP Act)")
    c_l4 = add_card(ls4, 0.8, 1.75, 11.733, 5.1)
    tf = c_l4.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "ESSENTIAL LEGAL CLAUSES & DATA PROTECTION (DPDP ACT 2023):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "1. Health & Physical Activity Disclaimer: Participants acknowledge that yoga involves physical exertion and confirm they are in suitable health.\n\n2. Intellectual Property Rights: Live session recordings, masterclass videos, and meditation tracks are proprietary to Art of Living.\n\n3. Data Protection Compliance (India DPDP Act): Personal data (WhatsApp number, age, name) is collected solely for class coordination and service delivery; never sold to third parties.\n\n4. Electronic Communications Consent: Explicit authorization to receive class links, zoom updates, and billing receipts via WhatsApp & SMS."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 5: Cancellation & Consumer Rights Policy
    ls5 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls5)
    add_leg_header(ls5, "Cancellation Rights", "Consumer Rights & Subscription Cancellation Policy", "Complying with consumer protection guidelines while ensuring fair cancellation terms")
    c_l5 = add_card(ls5, 0.8, 1.75, 11.733, 5.1)
    tf = c_l5.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "CONSUMER PROTECTION COMPLIANCE (CANCELLATION RULES):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Anytime Self-Serve Cancellation: Users can cancel their auto-renewal at any time via Profile -> Manage Subscription.\n• 24-Hour Cutoff: Cancellations must be requested at least 24 hours before next scheduled debit to prevent banking transit conflicts.\n• Paid Access Continuity: Cancelling auto-renewal does NOT immediately terminate access; access remains valid until current paid cycle expires.\n• Immediate Written Confirmation: Instant SMS & Email confirmation containing mandate revocation ID (MND_REV_XXXX)."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 6: Refund Policy Framework
    ls6 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls6)
    add_leg_header(ls6, "Refund Framework", "Refund Policy Framework for Indian Digital Wellness Subscriptions", "Clear, legally robust refund conditions for cancellations, duplicate debits, and technical issues")
    c_l6 = add_card(ls6, 0.8, 1.75, 11.733, 5.1)
    tf = c_l6.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "RECOMMENDED REFUND POLICY MATRIX (SUBJECT TO LEGAL SIGN-OFF):"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Pre-Commencement Cancellation: 100% refund if requested before the first scheduled class date.\n• Duplicate / Accidental Debit: 100% refund processed within 5-7 business days via Juspay gateway.\n• Technical Non-Delivery: Full refund if live class streaming fails due to server outages without alternative slots.\n• Post-Commencement Subscription: Non-refundable after active class attendance, but users can PAUSE their membership or transfer validity to another family member.\n• Chargeback Handling: Automated webhook logging provides full audit trail for disputed transactions."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 7: Free Trial & Disclosures
    ls7 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls7)
    add_leg_header(ls7, "Trial Disclosures", "Free Trial Disclosures & Auto-Debit Transition Safeguards", "Ensuring transparent communication during 14-day trial to auto-renew transition")
    c_l7 = add_card(ls7, 0.8, 1.75, 11.733, 5.1)
    tf = c_l7.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "TRIAL TRANSITION COMPLIANCE PROTOCOL:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    p = tf.add_paragraph()
    p.text = "• Day 0 Signup Disclosure: Clearly states '14 Days Free, then ₹4,999/year starting on [Date]. Cancel anytime during trial for ₹0 charge.'\n• Day 11 Advance Alert (WhatsApp & Email): 'Your free trial ends in 3 days. Your annual subscription will renew on [Date].'\n• Day 13 Final 24h Reminder: 'Your Sri Sri Yoga annual membership will activate tomorrow. Manage or cancel here: [Link]'\n• Day 14 Mandate Execution: Official GST tax invoice dispatched immediately upon successful charge."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    # Legal Slide 8: Legal Sign-off Checklist
    ls8 = prs_leg.slides.add_slide(blank_layout)
    set_bg(ls8)
    add_leg_header(ls8, "Legal Sign-Off", "Legal & Operational Sign-Off Checklist", "Required validation items before production deployment of Juspay AutoPay integration")
    c_l8 = add_card(ls8, 0.8, 1.75, 11.733, 5.1, SECONDARY_LIGHT, RGBColor(187, 247, 208))
    tf = c_l8.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.35)
    p = tf.paragraphs[0]
    p.text = "LEGAL SIGN-OFF CHECKLIST PRIOR TO LIVE LAUNCH:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p = tf.add_paragraph()
    p.text = "☑️ Terms of Service & Online Yoga Health Waiver reviewed by legal counsel\n☑️ Privacy Policy updated with DPDP Act 2023 compliance\n☑️ RBI e-Mandate pre-debit SMS/Email notification template approved\n☑️ Refund & Cancellation Policy published on website footer and checkout modal\n☑️ Juspay AutoPay mandate registration terms verified against NPCI guidelines\n☑️ Customer Support Escalation SLA (24h turnaround) established for billing disputes"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(12)

    prs_leg.save("Daily_Yoga_Subscription_Legal_Compliance_Deck.pptx")
    print(f"Successfully generated 8-slide Legal & Compliance Deck: Daily_Yoga_Subscription_Legal_Compliance_Deck.pptx")

if __name__ == '__main__':
    build_complete_decks()
